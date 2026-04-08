"""
运动员血液指标分析系统 - 终极修复版 v2.0
✅ 已修复：二级标题合并居中 + 颜色#C9DDE3 + 一级标题大字体+小间距

🎨 样式修改说明：
- 所有样式配置集中在第33-61行
- 颜色、字体、间距都可以在那里修改
- 详细说明请查看《样式修改指南.md》

🔥 版本标记：2.0 - 如果启动时看不到这个版本号，说明用的是旧文件！
"""

import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
import matplotlib
import numpy as np
from datetime import datetime
from io import BytesIO
from scipy.interpolate import make_interp_spline

# ========== 中文字体配置 ==========
import matplotlib.font_manager as fm
import os

font_path = os.path.join(os.path.dirname(__file__), 'fonts', 'SimHei.ttf')

if os.path.exists(font_path):
    fm.fontManager.addfont(font_path)
    plt.rcParams['font.sans-serif'] = ['SimHei']
    plt.rcParams['axes.unicode_minus'] = False
    print(f"✅ 成功加载中文字体：{font_path}")
else:
    print(f"❌ 字体文件不存在：{font_path}")
    plt.rcParams['font.sans-serif'] = ['DejaVu Sans']

from config import (
    MALE_REF_RANGES, FEMALE_REF_RANGES,
    COLUMN_NAME_MAPPING
)

# 趋势图默认指标
TREND_INDICATORS = ['睾酮', '皮质醇', '肌酸激酶', '血尿素', '血红蛋白', '铁蛋白', '白细胞', '网织红细胞百分比']


# ============================================================================
# 🎨 样式配置区域 - 在这里修改所有样式
# ============================================================================

# 【样式1】颜色配置
# 说明：修改这里可以改变所有颜色
# ⭐ 新配色方案：
# - 正常 → 浅灰色（白色看不见！）
# - 偏低/偏高 → 黄色
# - 严重偏低/严重偏高 → 红色
# - 良好/优秀 → 绿色
COLOR_SEVERE_LOW = '#FF6B6B'         # 严重偏低 - 红色
COLOR_LOW = '#FFD93D'                # 偏低 - 黄色
COLOR_NORMAL = '#F5F5F5'             # ⭐ 正常 - 浅灰色（原来白色看不见！）
COLOR_HIGH = '#FFD93D'               # 偏高 - 黄色（和偏低一样）
COLOR_SEVERE_HIGH = '#FF6B6B'        # 严重偏高 - 红色（和严重偏低一样）
COLOR_GOOD = '#6BCF7F'               # 良好 - 绿色
COLOR_EXCELLENT = '#6BCF7F'          # 优秀 - 绿色（和良好一样）

# ⭐【修改1】二级标题底色改为白色
COLOR_CATEGORY_HEADER = '#FFFFFF'    # 白色（分类标题背景）
COLOR_TABLE_HEADER = '#E8E8E8'       # 浅灰色（表头背景）- 改为浅灰，让黑色文字可见

COLOR_CHART_BG = '#FFFFFF'           # ⭐ 图表背景 - 全白
COLOR_MAIN = '#1f77b4'               # 主色调

# 【样式2】字体大小配置
# 说明：修改这里可以改变所有字体大小
FONTSIZE_MAIN_TITLE = 28    # 一级标题字体大小（增大）
FONTSIZE_HEADER = 18                 # 表头字体大小
FONTSIZE_CATEGORY = 16               # ⭐【修改2】分类标题字体（二级标题）- 原来是11
FONTSIZE_INDICATOR = 14             # ⭐【修改3】指标名称字体 - 原来是9
FONTSIZE_VALUE = 14                 # ⭐【修改3】数值字体 - 原来是10
FONTSIZE_STATUS = 14               # ⭐【修改3】状态字体 - 原来是8.5

# 【样式3】间距配置
# 说明：修改这里可以改变标题和表格的间距
TITLE_TABLE_SPACING = -0.5            # ⭐【修改4】一级标题和表格间距 - 原来是0.5，现在更小
TABLE_ROW_HEIGHT = 4               # ⭐【修改5】表格行高 - 从3增加到4，容纳多行标题

# ============================================================================
# 🔥 版本验证 - 启动时会在终端显示
# ============================================================================
print("=" * 60)
print("🚀 运动员血液指标分析系统 - v19.0 简化版")
print("=" * 60)
print(f"✅ 睾酮/皮质醇比值: 从Excel直接读取（不再计算）")
print(f"✅ 精准匹配 + 保留4位小数")
print(f"✅ 所有背景色统一")
print("=" * 60)
print("🎨 配色方案:")
print(f"   正常/无评价: {COLOR_NORMAL} (浅灰色)")
print(f"   偏低/偏高: {COLOR_LOW}/{COLOR_HIGH} (黄色)")
print(f"   严重: {COLOR_SEVERE_LOW}/{COLOR_SEVERE_HIGH} (红色)")
print(f"   优秀: {COLOR_GOOD}/{COLOR_EXCELLENT} (绿色)")
print("=" * 60)
# 测试颜色转换
from matplotlib.colors import to_rgba
test_color = to_rgba(COLOR_CATEGORY_HEADER)
print(f"🎨 二级标题颜色转换测试: {COLOR_CATEGORY_HEADER} → RGBA{test_color}")
print("=" * 60)

# ============================================================================
def check_password():
    def password_entered():
        if st.session_state["password"] == "blood2026":  # ← 改成你的密码
            st.session_state["password_correct"] = True
        else:
            st.session_state["password_correct"] = False

    if "password_correct" not in st.session_state:
        st.text_input("密码", type="password", on_change=password_entered, key="password")
        return False
    elif not st.session_state["password_correct"]:
        st.text_input("密码", type="password", on_change=password_entered, key="password")
        st.error("密码错误")
        return False
    return True


if not check_password():
    st.stop()

def parse_range_value(value_str):
    """
    解析范围值字符串

    支持格式：
    - "210-430" → (210, 430)
    - "< 210" → (None, 210)
    - "> 500" → (500, None)
    - "36.63" → (36.63, 36.63)
    - "6.0-20.0" → (6.0, 20.0)
    - "-" → (None, None)
    """
    if pd.isna(value_str) or str(value_str).strip() == '-' or str(value_str).strip() == '':
        return None, None

    value_str = str(value_str).strip()

    # 处理 "< X" 格式
    if value_str.startswith('<'):
        val = value_str.replace('<', '').strip()
        try:
            return None, float(val)
        except:
            return None, None

    # 处理 "> X" 格式
    if value_str.startswith('>'):
        val = value_str.replace('>', '').strip()
        try:
            return float(val), None
        except:
            return None, None

    # 处理 "X-Y" 格式
    if '-' in value_str:
        parts = value_str.split('-')
        if len(parts) == 2:
            try:
                return float(parts[0].strip()), float(parts[1].strip())
            except:
                return None, None

    # 处理单个数值
    try:
        val = float(value_str)
        return val, val
    except:
        return None, None


def load_reference_ranges_from_excel(file):
    """
    从上传的Excel文件加载参考范围

    返回：
    - male_ranges: 男性参考范围字典
    - female_ranges: 女性参考范围字典
    """
    try:
        # 读取参考范围sheet
        df = pd.read_excel(file, sheet_name='参考范围')

        male_ranges = {}
        female_ranges = {}
        common_ranges = {}

        # 遍历每一行
        for idx, row in df.iterrows():
            indicator = str(row['指标名称']).strip()
            gender = str(row['性别']).strip()

            # 解析五档范围
            severe_low_val = row['严重偏低 (<)']
            low_range = row['偏低 (范围)']
            normal_range = row['参考范围 (正常)']
            high_range = row['偏高 (范围)']
            severe_high_val = row['严重偏高 (>)']

            # 解析正常范围（这是最重要的）
            normal_low, normal_high = parse_range_value(normal_range)

            # 解析其他范围
            severe_low_lower, severe_low_upper = parse_range_value(severe_low_val)
            low_lower, low_upper = parse_range_value(low_range)
            high_lower, high_upper = parse_range_value(high_range)
            severe_high_lower, severe_high_upper = parse_range_value(severe_high_val)

            # 构建范围字典
            range_dict = {
                'severe_low_1': severe_low_lower if severe_low_lower is not None else severe_low_upper,
                'low_1': low_lower if low_lower is not None else None,
                'low_2': normal_low,  # 正常范围下限
                'high_2': normal_high,  # 正常范围上限
                'high_1': high_upper if high_upper is not None else None,
                'severe_high_1': severe_high_upper if severe_high_upper is not None else severe_high_lower,
            }

            # 根据性别分类
            if gender == '男':
                male_ranges[indicator] = range_dict
            elif gender == '女':
                female_ranges[indicator] = range_dict
            elif gender == '通用':
                common_ranges[indicator] = range_dict

        # 合并通用范围到男女范围
        for indicator, range_dict in common_ranges.items():
            if indicator not in male_ranges:
                male_ranges[indicator] = range_dict
            if indicator not in female_ranges:
                female_ranges[indicator] = range_dict

        return male_ranges, female_ranges

    except Exception as e:
        st.error(f"解析参考范围文件出错：{str(e)}")
        return {}, {}


# 设置中文字体
matplotlib.rcParams['font.sans-serif'] = ['Arial Unicode MS', 'SimHei', 'DejaVu Sans']
matplotlib.rcParams['axes.unicode_minus'] = False

# ========== 页面配置 ==========
st.set_page_config(
    page_title="运动员血液指标分析系统 - 增强版",
    page_icon="🏃",
    layout="wide",
    initial_sidebar_state="expanded"
)

# ========== 增强配置 ==========

# 主题配置 - 用于表格图
# 新的主题配置 - 7个生理系统分类
# 格式: '指标key': ('中文名', '英文名')

# 一级分类中英文对照
CATEGORY_NAMES = {
    '1_调控与指挥中心': ('调控与指挥中心（神经-内分泌系统）', 'Control and Command Center (Neuroendocrine System)'),
    '2_执行与代谢系统': ('执行与代谢系统（肌肉与能量状态）', 'Execution and Metabolic System (Muscle and Energy Status)'),
    '3_循环与运载系统': ('循环与运载系统（血液运载能力）', 'Circulation and Transport System (Blood Transport Capacity)'),
    '4_后勤保障与维护': ('后勤保障与维护（免疫与内环境）', 'Logistics Support and Maintenance (Immunity and Internal Environment)'),
}

THEME_CONFIG = {
    '1_调控与指挥中心': {
        '合成代谢\nAnabolism': {
            '睾酮': ('睾酮', 'Testosterone'),
            '游离睾酮': ('游离睾酮', 'Free Testosterone'),
        },
        '分解代谢\nCatabolism': {
            '皮质醇': ('皮质醇', 'Cortisol'),
        },
        '状态平衡\nStatus Balance': {
            '睾酮/皮质醇比值': ('睾酮/皮质醇比值', 'T/C Ratio'),
        }
    },

    '2_执行与代谢系统': {
        '结构完整性（硬件）\nStructural Integrity (Hardware)': {
            '肌酸激酶': ('肌酸激酶', 'Creatine Kinase'),
        },
        '能量储备与代谢（软件/燃料）\nEnergy Reserves and Metabolism (Software/Fuel)': {
            '血糖': ('血糖', 'Blood Glucose'),
            '血尿素': ('血尿素', 'Blood Urea'),
            '尿酸': ('尿酸', 'Uric Acid'),
        }
    },

    '3_循环与运载系统': {
        '输送载体（红细胞）\nTransport Carrier (Red Blood Cells)': {
            '血红蛋白': ('血红蛋白', 'Hemoglobin'),
            '红细胞': ('红细胞', 'RBC Count'),
            '红细胞压积': ('红细胞压积', 'Hematocrit'),
            '网织红细胞百分比': ('网织红细胞百分比', 'Reticulocyte %'),
            '平均红细胞容积': ('平均红细胞容积', 'MCV'),
        },
        '生化原料（造血储备）\nBiochemical Raw Materials (Hematopoietic Reserves)': {
            '铁蛋白': ('铁蛋白', 'Ferritin'),
            '维生素B12': ('维生素B12', 'Vitamin B12'),
            '维生素B6（PA）': ('维生素B6（PA）', 'Vitamin B6 (PA)'),  # ⭐ 修改
            '叶酸': ('叶酸', 'Folic Acid'),
        }
    },

    '4_后勤保障与维护': {
        '免疫防御（炎性监控）\nImmune Defense (Inflammatory Monitoring)': {
            '白细胞': ('白细胞', 'WBC Count'),
            '超敏C反应蛋白': ('超敏C反应蛋白', 'hs-CRP'),
            '触珠蛋白': ('触珠蛋白', 'Haptoglobin'),
        },
        '代谢辅酶（微量营养）\nMetabolic Coenzymes (Micronutrients)': {
            '维生素B1': ('维生素B1', 'Vitamin B1'),
            '维生素B2': ('维生素B2', 'Vitamin B2'),
            '维生素D2': ('维生素D2', 'Vitamin D2'),
            '维生素D3': ('维生素D3', 'Vitamin D3'),
            '维生素D': ('维生素D', 'Vitamin D'),
        },
        '内环境稳态（水盐平衡）\nInternal Environment Homeostasis (Water-Electrolyte Balance)': {
            '钾': ('钾', 'Potassium'),
            '钠': ('钠', 'Sodium'),
            '氯': ('氯', 'Chloride'),
            '渗透压': ('渗透压', 'Osmotic Pressure'),
            '血尿素/肌酐': ('血尿素/肌酐', 'BUN/Cr Ratio'),
        }
    },
}

# 雷达图配置
RADAR_FIELDS = ['游离睾酮', '皮质醇', '肌酸激酶', '血尿素', '血红蛋白', '铁蛋白', '白细胞', '网织红细胞百分比']  # ⭐ 改为游离睾酮
LOWER_IS_BETTER = ['肌酸激酶', '血尿素', '超敏C反应蛋白', '皮质醇']

# 颜色配置 - 五档评价配色
CATEGORY_NAMES = {
    '1_调控与指挥中心': ('调控与指挥中心（神经-内分泌系统）', 'Control and Command Center (Neuroendocrine System)'),
    '2_执行与代谢系统': ('执行与代谢系统（肌肉与能量状态）', 'Execution and Metabolic System (Muscle and Energy Status)'),
    '3_循环与运载系统': ('循环与运载系统（血液运载能力）', 'Circulation and Transport System (Blood Transport Capacity)'),
    '4_后勤保障与维护': ('后勤保障与维护（免疫与内环境）', 'Logistics Support and Maintenance (Immunity and Internal Environment)'),
}

THEME_CONFIG = {
    '1_调控与指挥中心': {
        '合成代谢\nAnabolism': {
            '睾酮': ('睾酮', 'Testosterone'),
            '游离睾酮': ('游离睾酮', 'Free Testosterone'),
        },
        '分解代谢\nCatabolism': {
            '皮质醇': ('皮质醇', 'Cortisol'),
        },
        '状态平衡\nStatus Balance': {
            '睾酮/皮质醇比值': ('睾酮/皮质醇比值', 'T/C Ratio'),
        }
    },

    '2_执行与代谢系统': {
        '结构完整性（硬件）\nStructural Integrity (Hardware)': {
            '肌酸激酶': ('肌酸激酶', 'Creatine Kinase'),
        },
        '能量储备与代谢（软件/燃料）\nEnergy Reserves and Metabolism (Software/Fuel)': {
            '血糖': ('血糖', 'Blood Glucose'),
            '血尿素': ('血尿素', 'Blood Urea'),
            '尿酸': ('尿酸', 'Uric Acid'),
        }
    },

    '3_循环与运载系统': {
        '输送载体（红细胞）\nTransport Carrier (Red Blood Cells)': {
            '血红蛋白': ('血红蛋白', 'Hemoglobin'),
            '红细胞': ('红细胞', 'RBC Count'),
            '红细胞压积': ('红细胞压积', 'Hematocrit'),
            '网织红细胞百分比': ('网织红细胞百分比', 'Reticulocyte %'),
            '平均红细胞容积': ('平均红细胞容积', 'MCV'),
        },
        '生化原料（造血储备）\nBiochemical Raw Materials (Hematopoietic Reserves)': {
            '铁蛋白': ('铁蛋白', 'Ferritin'),
            '维生素B12': ('维生素B12', 'Vitamin B12'),
            '维生素B6（PA）': ('维生素B6（PA）', 'Vitamin B6 (PA)'),  # ⭐ 修改
            '叶酸': ('叶酸', 'Folic Acid'),
        }
    },

    '4_后勤保障与维护': {
        '免疫防御（炎性监控）\nImmune Defense (Inflammatory Monitoring)': {
            '白细胞': ('白细胞', 'WBC Count'),
            '超敏C反应蛋白': ('超敏C反应蛋白', 'hs-CRP'),
            '触珠蛋白': ('触珠蛋白', 'Haptoglobin'),
        },
        '代谢辅酶（微量营养）\nMetabolic Coenzymes (Micronutrients)': {
            '维生素B1': ('维生素B1', 'Vitamin B1'),
            '维生素B2': ('维生素B2', 'Vitamin B2'),
            '维生素D2': ('维生素D2', 'Vitamin D2'),
            '维生素D3': ('维生素D3', 'Vitamin D3'),
            '维生素D': ('维生素D', 'Vitamin D'),
        },
        '内环境稳态（水盐平衡）\nInternal Environment Homeostasis (Water-Electrolyte Balance)': {
            '钾': ('钾', 'Potassium'),
            '钠': ('钠', 'Sodium'),
            '氯': ('氯', 'Chloride'),
            '渗透压': ('渗透压', 'Osmotic Pressure'),
            '血尿素/肌酐': ('血尿素/肌酐', 'BUN/Cr Ratio'),
        }
    },
}

# 雷达图配置
RADAR_FIELDS = ['游离睾酮', '皮质醇', '肌酸激酶', '血尿素', '血红蛋白', '铁蛋白', '白细胞', '网织红细胞百分比']  # ⭐ 改为游离睾酮
LOWER_IS_BETTER = ['肌酸激酶', '血尿素', '超敏C反应蛋白', '皮质醇']

# 颜色配置 - 五档评价配色（新配色方案）
COLOR_SEVERE_LOW = '#FF6B6B'     # 红色（严重偏低）
COLOR_LOW = '#FFD93D'            # 黄色（偏低）
COLOR_NORMAL = '#F5F5F5'         # 浅灰色（正常）- 原来白色看不见！
COLOR_HIGH = '#FFD93D'           # 黄色（偏高）
COLOR_SEVERE_HIGH = '#FF6B6B'    # 红色（严重偏高）
COLOR_GOOD = '#6BCF7F'           # 绿色（良好）
COLOR_EXCELLENT = '#6BCF7F'      # 绿色（优秀）
COLOR_CATEGORY_HEADER = '#FFFFFF'  # 白色（分类标题）
COLOR_CHART_BG = '#F8F9FA'       # 极浅灰（图表背景）
COLOR_MAIN = '#1f77b4'          # 主色调

# 雷达图样式
RADAR_STYLES = [
    {'color': '#8BC1E9', 'linewidth': 2, 'linestyle': ':'},   # 第1次 - 浅天蓝
    {'color': '#E89A9D', 'linewidth': 2, 'linestyle': '-.'},  # 第2次 - 浅柔红
    {'color': '#5C7CFA', 'linewidth': 2.5, 'linestyle': '--'}, # 第3次 - 靛蓝
    {'color': '#D05A5E', 'linewidth': 3, 'linestyle': '-'},   # 第4次（最新）- 深砖红
]

@st.cache_data(show_spinner=False)
def load_data_multisheet(file_path_or_buffer):
    """
    从多个sheet加载数据并合并 - 增强版（自动识别格式，无需固定sheet名）
    支持：任意sheet名称 / 任意表头行位置 / 有无10行间隔 / 双层表头
    """
    from io import BytesIO

    # ★ 关键修复：Streamlit 上传的文件对象只能读一次。
    #   先把全部内容读入内存（BytesIO），后续所有 pd.read_excel / ExcelFile
    #   都从同一个可重复 seek 的 BytesIO 读取，不会卡死。
    try:
        if hasattr(file_path_or_buffer, 'read'):
            raw_bytes = file_path_or_buffer.read()
        else:
            with open(file_path_or_buffer, 'rb') as fh:
                raw_bytes = fh.read()
        file_buf = BytesIO(raw_bytes)
    except Exception as e:
        st.error(f"❌ 读取文件失败：{e}")
        return None

    def is_athlete_name(s):
        """判断字符串是否为有效运动员姓名（2-6位中文）"""
        s = str(s).strip()
        return (2 <= len(s) <= 6 and
                all('\u4e00' <= c <= '\u9fa5' for c in s))

    # ── Step 1: 探测所有 Sheet ──────────────────────────────────────────────
    # 用 openpyxl read_only 模式快速取 sheet 名（比 pd.ExcelFile 快 3-5x）
    try:
        from openpyxl import load_workbook as _opxl_load
        file_buf.seek(0)
        _wb = _opxl_load(file_buf, read_only=True, data_only=True)
        all_sheets = _wb.sheetnames
        _wb.close()
    except Exception:
        try:
            file_buf.seek(0)
            xl = pd.ExcelFile(file_buf)
            all_sheets = xl.sheet_names
        except Exception as e:
            st.error(f"❌ 无法打开Excel文件：{e}")
            st.info("💡 支持格式：.xlsx / .xls；请确认文件未加密、未损坏")
            return None

    if not all_sheets:
        st.error("❌ Excel文件中没有找到任何Sheet")
        return None

    st.info(f"📋 发现 {len(all_sheets)} 个Sheet：{' | '.join(all_sheets)}")

    # ── Step 2: 找主数据 Sheet ─────────────────────────────────────────────
    MONTHLY_KEYWORDS = ['月周', '月测', '主要', '月', 'monthly', 'main', 'data']
    monthly_sheet = None

    # 精确匹配
    for sn in all_sheets:
        if sn.strip() == '月周测试指标':
            monthly_sheet = sn
            break

    # 关键词匹配
    if not monthly_sheet:
        for sn in all_sheets:
            if any(k in sn for k in MONTHLY_KEYWORDS):
                monthly_sheet = sn
                break

    # 找含"姓名"的 Sheet
    if not monthly_sheet:
        for sn in all_sheets:
            try:
                file_buf.seek(0)
                df_peek = pd.read_excel(file_buf, sheet_name=sn,
                                        nrows=25, header=None, dtype=str)
                if df_peek.astype(str).apply(
                    lambda col: col.str.contains('姓名', na=False)
                ).any().any():
                    monthly_sheet = sn
                    break
            except Exception:
                continue

    # 兜底：第一个 Sheet
    if not monthly_sheet:
        monthly_sheet = all_sheets[0]
        st.warning(f"⚠️ 未找到标准Sheet名，自动使用：**{monthly_sheet}**")
    else:
        st.write(f"✓ 主数据 Sheet：**{monthly_sheet}**")

    # ── Step 3: 读取原始数据 ──────────────────────────────────────────────
    # 两步读取：先25行定位表头，再精准读取完整数据（大文件速度提升 3-5x）
    try:
        file_buf.seek(0)
        df_peek = pd.read_excel(file_buf, sheet_name=monthly_sheet,
                                nrows=25, header=None, dtype=str, engine='openpyxl')
    except Exception as e:
        st.error(f"❌ 读取 Sheet [{monthly_sheet}] 失败：{e}")
        return None
    df_peek = df_peek.fillna('')

    _h_tmp, _nc_tmp = None, 0
    _NV = {'姓名','运动员姓名','运动员','Name','name'}
    for _i in range(len(df_peek)):
        for _j in range(len(df_peek.columns)):
            if str(df_peek.iloc[_i,_j]).strip() in _NV:
                _h_tmp = _i; _nc_tmp = _j; break
        if _h_tmp is not None: break

    _skip = list(range(1, _h_tmp)) if (_h_tmp and _h_tmp > 1) else []
    try:
        file_buf.seek(0)
        df_raw = pd.read_excel(file_buf, sheet_name=monthly_sheet,
                               header=None, dtype=str,
                               skiprows=_skip, engine='openpyxl')
    except Exception as e:
        st.error(f"❌ 读取完整数据失败：{e}")
        return None
    df_raw = df_raw.fillna('')

    # ── Step 4: 找表头行 ──────────────────────────────────────────────────
    header_row = None
    name_col_idx = 0
    NAME_VARIANTS = {'姓名', '运动员姓名', '运动员', 'Name', 'name'}

    for i in range(min(30, len(df_raw))):
        for j in range(len(df_raw.columns)):
            if str(df_raw.iloc[i, j]).strip() in NAME_VARIANTS:
                header_row = i
                name_col_idx = j
                break
        if header_row is not None:
            break

    if header_row is None:
        st.error("❌ 未在前30行找到【姓名】列")
        with st.expander("📊 前15行内容（点击展开帮助诊断）", expanded=True):
            st.dataframe(df_raw.head(15))
        st.info("💡 请确认数据表中有【姓名】作为列标题，且位于前30行内")
        return None

    st.write(f"✓ 表头行：第 **{header_row + 1}** 行，姓名列：第 **{name_col_idx + 1}** 列")

    # ── Step 5: 检测双层表头 ─────────────────────────────────────────────
    CATEGORY_NAMES = {'维生素', '电解质', '甲功', '肝功', '血脂', '甲状腺', '血脂四项', '糖类'}
    headers = list(df_raw.iloc[header_row])
    is_double_hdr = any(str(h).strip() in CATEGORY_NAMES for h in headers)

    if is_double_hdr and header_row + 1 < len(df_raw):
        next_row = list(df_raw.iloc[header_row + 1])
        headers = [
            str(next_row[i]).strip()
            if str(next_row[i]).strip() and not str(next_row[i]).startswith('Unnamed')
            else str(headers[i]).strip()
            for i in range(len(headers))
        ]
        header_row += 1
        st.write("✓ 检测到双层表头，已合并")

    # ── Step 6: 跳过非运动员行，找数据起始行 ────────────────────────────
    data_start = header_row + 1
    for i in range(header_row + 1, min(header_row + 16, len(df_raw))):
        cell = str(df_raw.iloc[i, name_col_idx]).strip()
        if is_athlete_name(cell):
            data_start = i
            break

    st.write(f"✓ 数据起始行：第 **{data_start + 1}** 行")

    # ── Step 7: 构建运动员数据 DataFrame ───────────────────────────────
    data_rows = df_raw.iloc[data_start:].copy()
    data_rows.columns = range(len(data_rows.columns))

    # 过滤有效运动员行
    valid_mask = data_rows.iloc[:, name_col_idx].apply(
        lambda x: is_athlete_name(str(x))
    )
    data_rows = data_rows[valid_mask].copy()

    if len(data_rows) == 0:
        st.error("❌ 未在数据行中找到有效运动员姓名（需为2-6位纯汉字）")
        with st.expander("📊 表头后内容预览"):
            st.dataframe(df_raw.iloc[header_row: header_row + 10])
        return None

    # 赋列名（处理列数不匹配）
    min_cols = min(len(headers), len(data_rows.columns))
    data_rows = data_rows.iloc[:, :min_cols].copy()
    data_rows.columns = headers[:min_cols]

    # 数值转换
    NON_NUMERIC = {'姓名', '运动员姓名', '运动员', '测试日期', '日期', '性别', '备注', 'Name'}
    for col in data_rows.columns:
        if str(col) not in NON_NUMERIC:
            data_rows[col] = pd.to_numeric(data_rows[col], errors='ignore')

    # 列名唯一化
    seen, new_cols = {}, []
    for c in data_rows.columns:
        cs = str(c)
        if cs in seen:
            seen[cs] += 1
            new_cols.append(f"{cs}#{seen[cs]}")
        else:
            seen[cs] = 0
            new_cols.append(cs)
    data_rows.columns = new_cols
    df_monthly = data_rows.reset_index(drop=True)
    st.write(f"✓ 月周测试：**{len(df_monthly)}** 行，**{len(df_monthly.columns)}** 列")

    # ── Step 8: 合并其他 Sheet ──────────────────────────────────────────
    QUARTER_KW = ['季度', '季', '维生素', '电解质', '微量', 'quarterly']
    YEAR_KW    = ['年度', '年', '甲功', '肝功', '血脂', 'annual', 'yearly']
    OTHER_KW   = ['其他', '额外', 'other', '触珠']

    name_col_final = '姓名' if '姓名' in df_monthly.columns else df_monthly.columns[name_col_idx]
    date_col_final = next((c for c in df_monthly.columns
                           if '日期' in str(c) or str(c) in {'Date', 'DateStr'}), None)
    df_merged = df_monthly.copy()

    if date_col_final:
        df_merged['_merge_key'] = (
            df_merged[name_col_final].astype(str) + '_' +
            df_merged[date_col_final].astype(str)
        )

    def find_sheet_by_kw(kws, exclude):
        for sn in all_sheets:
            if sn == exclude:
                continue
            if any(k in sn for k in kws):
                return sn
        return None

    for kws, label in [(QUARTER_KW, '季度'), (YEAR_KW, '年度'), (OTHER_KW, '其他')]:
        sn = find_sheet_by_kw(kws, monthly_sheet)
        if not sn:
            continue
        try:
            file_buf.seek(0)
            df_add_raw = pd.read_excel(file_buf, sheet_name=sn, header=[0, 1])
            df_add = flatten_multiindex_columns(df_add_raw, label)
            df_merged = merge_sheet_data(df_merged, df_add,
                                         name_col_final, date_col_final, label)
            st.write(f"✓ 已合并 [{sn}] Sheet（{label}数据）")
        except Exception as e:
            st.warning(f"⚠️ [{sn}] 合并失败（{e}），跳过")

    if '_merge_key' in df_merged.columns:
        df_merged.drop('_merge_key', axis=1, inplace=True)

    # ── Step 9: 列名标准化 ──────────────────────────────────────────────
    df_merged = df_merged.rename(columns=COLUMN_NAME_MAPPING)

    # 处理重复列（保留第一列）
    if df_merged.columns.duplicated().any():
        seen2 = {}
        keep = []
        for i, c in enumerate(df_merged.columns):
            if c not in seen2:
                seen2[c] = True
                keep.append(i)
        df_merged = df_merged.iloc[:, keep]

    athletes_found = sorted(df_merged[name_col_final].dropna().unique())
    st.success(f"✅ 数据加载完成：{len(df_merged)} 条记录，{len(athletes_found)} 名运动员")
    st.write(f"🏃 运动员列表：{' / '.join(map(str, athletes_found))}")

    return df_merged




def flatten_multiindex_columns(df, sheet_name):
    """
    将双层MultiIndex列名展平为单层
    优先使用具体的指标名称
    """
    # 定义分类名称（这些应该跳过，使用level1）
    category_names = ['维生素', '电解质', '甲功', '肝功', '血脂四项', '糖类指标', '性别']
    
    # 定义基础信息列名（这些应该保留level0）
    basic_info_cols = ['项目', '编号', '姓名', '出生年月日', '身高', '体重', '测试日期', '备注']
    
    new_columns = []
    
    for col in df.columns:
        if isinstance(col, tuple):
            level0, level1 = col[0], col[1]
            
            # 判断逻辑：
            # 1. 如果level0是基础信息列，使用level0
            # 2. 如果level0是分类名称，使用level1（指标名）
            # 3. 如果level1是Unnamed，使用level0
            # 4. 否则优先使用level1
            
            if str(level0) in basic_info_cols:
            # ⭐ 修改：基础信息列，只有当level1无效时才使用level0
                if pd.isna(level1) or str(level1).startswith('Unnamed'):
                    col_name = str(level0)
                else:
        # 如果level1有效，使用level1（实际指标名）
                    col_name = str(level1)
            elif str(level0) in category_names:
                # 分类名称：使用level1（实际指标名）
                if not (pd.isna(level1) or str(level1).startswith('Unnamed')):
                    col_name = str(level1)
                else:
                    col_name = f'Unnamed_{len(new_columns)}'
            elif pd.isna(level1) or str(level1).startswith('Unnamed'):
                # level1无效：使用level0
                if not (pd.isna(level0) or str(level0).startswith('Unnamed')):
                    col_name = str(level0)
                else:
                    col_name = f'Unnamed_{len(new_columns)}'
            else:
                # 其他情况：优先使用level1
                col_name = str(level1)
        else:
            col_name = str(col)
        
        new_columns.append(col_name)
    
    df.columns = new_columns
    
    # 🔧 确保列名唯一
    seen = {}
    unique_columns = []
    for col in df.columns:
        if col in seen:
            seen[col] += 1
            unique_columns.append(f"{col}_{seen[col]}")
        else:
            seen[col] = 0
            unique_columns.append(col)
    
    df.columns = unique_columns
    
    return df


def merge_all_sheets(df_monthly, df_quarterly, df_yearly, df_other):
    """
    合并所有sheet的数据
    使用姓名和测试日期作为合并键
    """
    # 以月周测试数据为基础
    df_result = df_monthly.copy()
    
    # 确定合并键 - 尝试多种可能的列名
    name_col_monthly = None
    for col_name in ['姓名', 'Name', 'Name_final']:
        if col_name in df_result.columns:
            name_col_monthly = col_name
            break
    
    date_col_monthly = None
    for col_name in ['测试日期', 'Date', 'Date_auto']:
        if col_name in df_result.columns:
            date_col_monthly = col_name
            break
    
    if not name_col_monthly or not date_col_monthly:
        st.warning("⚠ 无法找到姓名或日期列，仅使用月周测试数据")
        return df_result
    
    # 创建合并键
    df_result['_merge_key'] = (
        df_result[name_col_monthly].astype(str) + '_' + 
        df_result[date_col_monthly].astype(str)
    )
    
    # 合并季度测试数据
    if df_quarterly is not None:
        df_result = merge_sheet_data(df_result, df_quarterly, name_col_monthly, date_col_monthly, '季度测试')
    
    # 合并年度测试数据
    if df_yearly is not None:
        df_result = merge_sheet_data(df_result, df_yearly, name_col_monthly, date_col_monthly, '年度测试')
    
    # 合并其他数据
    if df_other is not None:
        df_result = merge_sheet_data(df_result, df_other, name_col_monthly, date_col_monthly, '其他')
    
    # 删除临时合并键
    if '_merge_key' in df_result.columns:
        df_result = df_result.drop('_merge_key', axis=1)
    
    # 🔧 新增：应用列名映射，标准化列名
    df_result = df_result.rename(columns=COLUMN_NAME_MAPPING)
    st.write(f"   ✓ 列名标准化完成")
    
    # 🔧 检查并处理映射后的重复列名
    if df_result.columns.duplicated().any():
        st.warning("   ⚠ 映射后发现重复列名，正在处理...")
        
        # 对于重复列名，保留第一列，删除后续列
        seen = {}
        cols_to_keep = []
        for i, col in enumerate(df_result.columns):
            if col not in seen:
                seen[col] = i
                cols_to_keep.append(i)
        
        df_result = df_result.iloc[:, cols_to_keep]
        st.write(f"   ✓ 重复列名已处理，保留了{len(cols_to_keep)}列")
    
    return df_result


def merge_sheet_data(df_main, df_add, name_col, date_col, sheet_name):
    """
    将额外sheet的数据合并到主数据框
    """
    try:
        # 🔧 新增：检查df_add是否有重复列名
        if df_add.columns.duplicated().any():
            st.warning(f"   ⚠ {sheet_name}：发现重复列名，正在修复...")
            # 修复重复列名
            seen = {}
            unique_columns = []
            for col in df_add.columns:
                if col in seen:
                    seen[col] += 1
                    unique_columns.append(f"{col}_{seen[col]}")
                else:
                    seen[col] = 0
                    unique_columns.append(col)
            df_add.columns = unique_columns
            st.write(f"   ✓ 列名已唯一化")
        
        # 在额外sheet中找到对应的姓名和日期列
        name_col_add = None
        for col_name in ['姓名', 'Name', 'Name_final']:
            if col_name in df_add.columns:
                name_col_add = col_name
                break
        
        date_col_add = None
        for col_name in ['测试日期', 'Date', 'Date_auto']:
            if col_name in df_add.columns:
                date_col_add = col_name
                break
        
        if not name_col_add:
            st.warning(f"   ⚠ {sheet_name}：无法找到姓名列，跳过合并")
            return df_main
            
        if not date_col_add:
            st.warning(f"   ⚠ {sheet_name}：无法找到日期列，跳过合并")
            return df_main
        
        # 创建合并键（使用.copy()避免SettingWithCopyWarning）
        df_add = df_add.copy()
        df_add['_merge_key'] = (
            df_add[name_col_add].astype(str) + '_' + 
            df_add[date_col_add].astype(str)
        )
        
        # 选择要合并的指标列（排除基本信息列）
        exclude_cols = [
            '项目', '编号', '姓名', '性别', '出生年月日', '身高', '体重', '测试日期', 
            'Name', 'Name_final', 'Date', 'Date_auto', '_merge_key',
            '教练', '训练地点', '测试单位', '测试阶段', '重点运动员', '专项'
        ]
        
        indicator_cols = []
        for col in df_add.columns:
            if col in exclude_cols:
                continue
            if str(col).startswith('Unnamed'):
                continue
            if pd.isna(col):
                continue
            indicator_cols.append(col)
        
        if len(indicator_cols) == 0:
            st.warning(f"   ⚠ {sheet_name}：没有找到指标列，跳过合并")
            return df_main
        
        # 只保留指标列和合并键
        df_add_indicators = df_add[['_merge_key'] + indicator_cols].copy()
        
        # 🔧 新增：检查df_add_indicators是否有重复列名
        if df_add_indicators.columns.duplicated().any():
            st.warning(f"   ⚠ {sheet_name}：指标列有重复，正在去重...")
            # 再次确保唯一性
            seen = {}
            unique_columns = []
            for col in df_add_indicators.columns:
                if col in seen:
                    seen[col] += 1
                    unique_columns.append(f"{col}_{seen[col]}")
                else:
                    seen[col] = 0
                    unique_columns.append(col)
            df_add_indicators.columns = unique_columns
        
        # 合并数据
        df_merged = df_main.merge(
            df_add_indicators,
            on='_merge_key',
            how='left',
            suffixes=('', f'_from_{sheet_name}')
        )
        
        # 🔧 新增：合并后再次检查重复列名
        if df_merged.columns.duplicated().any():
            st.warning(f"   ⚠ 合并后发现重复列名，正在修复...")
            seen = {}
            unique_columns = []
            for col in df_merged.columns:
                if col in seen:
                    seen[col] += 1
                    unique_columns.append(f"{col}_{seen[col]}")
                else:
                    seen[col] = 0
                    unique_columns.append(col)
            df_merged.columns = unique_columns
        
        st.write(f"   ✓ {sheet_name}合并：添加了 {len(indicator_cols)} 个指标")
        
        return df_merged
        
    except Exception as e:
        st.warning(f"   ⚠ {sheet_name}合并失败：{e}")
        import traceback
        st.write(traceback.format_exc())
        return df_main



def clean_data_final(df):
    """数据清洗函数"""
    if df is None:
        return None

    st.info("🧹 开始清洗数据...")

    # 删除空行
    df = df.dropna(how='all')
    df = df.reset_index(drop=True)

    # 处理姓名列
    if '姓名' in df.columns:
        name_cols = [col for col in df.columns if col.startswith('Name')]
        if not name_cols:
            df['Name'] = df['姓名']
        else:
            df['Name_final'] = df['姓名']

    # 处理日期列
    possible_date_cols = ['测试日期', '日期', '开始日期']
    date_col_found = False

    for col in possible_date_cols:
        if col in df.columns:
            try:
                date_cols = [c for c in df.columns if c.startswith('Date')]

                if not date_cols:
                    if pd.api.types.is_datetime64_any_dtype(df[col]):
                        df['Date'] = df[col]
                    else:
                        df['Date'] = pd.to_datetime(df[col], errors='coerce')

                    df['DateStr'] = df['Date'].dt.strftime('%Y-%m-%d')
                    date_col_found = True
                else:
                    date_col_found = True

                break

            except Exception as e:
                continue

    if not date_col_found:
        df['Date_auto'] = pd.date_range(start='2024-01-01', periods=len(df), freq='D')
        df['DateStr'] = df['Date_auto'].dt.strftime('%Y-%m-%d')

    # 最终清理
    df = df.dropna(how='all')
    df = df.reset_index(drop=True)

    st.success(f"✅ 清洗完成：保留 {len(df)} 行有效数据")

    return df

# ========== 辅助函数 ==========

def get_indicator_status(indicator, value, ref_ranges, gender=None):
    """判断指标状态（五档）- 完全修复版
    
    参数:
    - indicator: 指标名称
    - value: 指标值
    - ref_ranges: 参考范围字典
    - gender: 性别（'男'或'女'），可选，用于铁蛋白特殊评价
    """
    # 先检查是否为NaN
    if indicator not in ref_ranges or pd.isna(value):
        return '-', COLOR_NORMAL, 'N/A'  # ⭐ 改为COLOR_NORMAL
    
    # 🔧 修复1：转换value为数值类型
    try:
        if isinstance(value, str):
            value = value.strip()
            if value == '' or value == '-' or value.lower() == 'nan':
                return '-', COLOR_NORMAL, 'N/A'  # ⭐ 改为COLOR_NORMAL
            value = float(value)
        elif not isinstance(value, (int, float)):
            value = float(value)
    except (ValueError, TypeError):
        return '-', COLOR_NORMAL, 'N/A'  # ⭐ 改为COLOR_NORMAL

    ranges = ref_ranges[indicator]
    
    # ⭐ 新增：检查是否为空字典（不评价的指标）
    if not ranges or len(ranges) == 0:
        return '-', COLOR_NORMAL, 'N/A'
    
    # 🔧 修复2：确保参考范围值也是数值类型
    try:
        low_1 = ranges.get('low_1')
        low_2 = ranges.get('low_2')
        high_2 = ranges.get('high_2')
        high_1 = ranges.get('high_1')
        
        # 转换参考范围值为浮点数
        if low_1 is not None and not isinstance(low_1, (int, float)):
            low_1 = float(low_1) if not pd.isna(low_1) else None
        if low_2 is not None and not isinstance(low_2, (int, float)):
            low_2 = float(low_2) if not pd.isna(low_2) else None
        if high_2 is not None and not isinstance(high_2, (int, float)):
            high_2 = float(high_2) if not pd.isna(high_2) else None
        if high_1 is not None and not isinstance(high_1, (int, float)):
            high_1 = float(high_1) if not pd.isna(high_1) else None
    except (ValueError, TypeError):
        return '-', COLOR_NORMAL, 'N/A'  # ⭐ 改为COLOR_NORMAL

    # 高优指标列表（高于正常范围是好事）
    high_is_better_indicators = ['铁蛋白', '血红蛋白', '睾酮', '游离睾酮']
    
    # ⭐ 新增：偏高不评价的指标列表（偏高时返回"正常"）
    no_high_evaluation_indicators = ['维生素B1', '维生素B2','维生素B12']
    
    # ⭐ 新增：铁蛋白特殊处理 - 过高需要注意
    if indicator == '铁蛋白' and gender:
        if (gender == '男' and value > 300) or (gender == '女' and value > 200):
            return '需注意', '#FFA500', 'attention_needed'  # 橙色

    # 🔧 修复3：判断状态时添加异常保护
    try:
        if pd.notna(low_1) and value < low_1:
            return '严重偏低', COLOR_SEVERE_LOW, 'severe_low'
        elif pd.notna(low_2) and value < low_2:
            return '偏低', COLOR_LOW, 'low'
        elif pd.notna(high_1) and value > high_1:
            # ⭐ 修改：如果是不评价偏高的指标，返回"-"
            if indicator in no_high_evaluation_indicators:
                return '-', COLOR_NORMAL, 'N/A'
            elif indicator in high_is_better_indicators:
                return '优秀', COLOR_EXCELLENT, 'excellent'  # 使用绿色
            else:
                return '严重偏高', COLOR_SEVERE_HIGH, 'severe_high'
        elif pd.notna(high_2) and value > high_2:
            # ⭐ 修改：如果是不评价偏高的指标，返回"-"
            if indicator in no_high_evaluation_indicators:
                return '-', COLOR_NORMAL, 'N/A'
            elif indicator in high_is_better_indicators:
                return '良好', COLOR_GOOD, 'good'  # 使用绿色
            else:
                return '偏高', COLOR_HIGH, 'high'
        else:
            return '正常', COLOR_NORMAL, 'normal'
    except (TypeError, ValueError):
        return '-', COLOR_NORMAL, 'N/A'  # ⭐ 改为COLOR_NORMAL


def format_number(val):
    """智能格式化数值，保留完整小数位但去除尾部0"""
    if pd.isna(val):
        return "—"
    try:
        val = float(val)
        # 如果是整数，直接显示整数
        if val == int(val):
            return f"{int(val)}"
        # 否则保留原始精度，但去除尾部0
        # 先格式化为字符串，保留足够精度
        formatted = f"{val:.10f}".rstrip('0').rstrip('.')
        return formatted
    except (ValueError, TypeError):
        return "—"


# 指标别名映射（用于处理常见的名称差异）
INDICATOR_ALIASES = {
    # 红细胞指标
    '平均红细胞血红蛋白浓度': ['平均红细胞血红浓度', 'MCHC', '平均血红蛋白浓度'],
    '平均红细胞血红蛋白': ['平均红细胞血红蛋白量', 'MCH'],
    '平均红细胞体积': ['平均红细胞容积', 'MCV'],
    '平均红细胞容积': ['平均红细胞体积', 'MCV'],
    '平均血红蛋白浓度': ['平均红细胞血红蛋白浓度', 'MCHC'],
    '网织红细胞百分比': ['网织红细胞', 'retic', 'Retic'],
    
    # 炎症指标
    '超敏C反应蛋白': ['C反应蛋白', 'CRP', 'hsCRP', 'hs-CRP'],
    
    # 维生素指标（季度测试）
    '维生素B1': ['VB1', 'VitB1'],
    '维生素B2': ['VB2', 'VitB2'],
    '维生素B6（PA）': ['VB6', 'VitB6', 'VitB6(PA)', 'B6'],  # ⭐ 修改
    '维生素B6（PLP）': ['vitB6（PLP）', 'VitB6(PLP)', 'B6(PLP)'],  # ⭐ 新增
    '维生素B12': ['VB12', 'VitB12'],
    '叶酸': ['FOL', '维生素B9'],
    '维生素D3': ['VD3', 'VD3(25-OH)', 'VD-(25-OH)'],
    
    # 电解质（季度测试）
    '钾': ['K'],
    '钠': ['Na'],
    '氯': ['Cl'],
    '钙': ['Ca'],
    '镁': ['Mg'],
    
    # 甲状腺功能（年度测试）
    '总甲状腺素': ['T4', 'TT4'],
    '总三碘甲状腺原氨酸': ['T3', 'TT3'],
    '游离三碘甲状原氨酸': ['FT3', '游离T3'],
    '游离甲状腺素': ['FT4', '游离T4'],
    '超敏促甲状腺素': ['TSH', 'hs-TSH', '促甲状腺激素'],
    
    # 肝功能（年度测试）
    '丙氨酸氨基转移酶': ['ALT', '谷丙转氨酶', '丙氨酸基转移酶'],
    '天冬氨酸氨基转移酶': ['AST', '谷草转氨酶'],
    '碱性磷酸酶': ['ALP'],
    'γ-谷氨酰基转移酶': ['GGT', 'γ-GT', 'γ-谷氨酰转移酶'],
    '总胆红素': ['TBIL', 'TB'],
    '直接胆红素': ['DBIL', 'DB'],
    '间接胆红素': ['IBIL', 'IB'],
    '总蛋白': ['TP'],
    '白蛋白': ['ALB', 'Alb'],
    
    # 血脂（年度测试）
    '甘油三酯': ['TG', 'TAG'],
    '高密度脂蛋白': ['HDL', 'HDL-C'],
    '总胆固醇': ['TC', 'CHOL'],
    '低密度脂蛋白': ['LDL', 'LDL-C'],
}

def find_indicator_column(df, indicator):
    """智能查找指标列（支持带#的列名、模糊匹配、别名匹配）"""

    # ⭐ 特殊处理：重要指标优先精确匹配（避免匹配到.1后缀的重复列）
    PRIORITY_INDICATORS = ['睾酮', '游离睾酮', '皮质醇', '睾酮/皮质醇比值']
    if indicator in PRIORITY_INDICATORS:
        # 优先精确匹配
        if indicator in df.columns:
            return indicator
        # 如果没有精确匹配，再尝试带#后缀的
        for col in df.columns:
            col_str = str(col)
            if col_str.startswith(indicator):
                suffix = col_str[len(indicator):]
                if suffix.startswith('#'):  # 只允许#后缀，不允许.数字
                    return col

    # 方法1：精确匹配
    if indicator in df.columns:
        return indicator

    # 方法2：别名匹配
    # 先查找是否有直接的别名定义
    if indicator in INDICATOR_ALIASES:
        for alias in INDICATOR_ALIASES[indicator]:
            if alias in df.columns:
                return alias
            # 也尝试前缀匹配别名
            possible_cols = [col for col in df.columns if str(col).startswith(alias)]
            if possible_cols:
                return possible_cols[0]

    # 反向查找：indicator是否是某个别名
    for main_name, aliases in INDICATOR_ALIASES.items():
        if indicator in aliases:
            # 尝试匹配主名称
            if main_name in df.columns:
                return main_name
            possible_cols = [col for col in df.columns if str(col).startswith(main_name)]
            if possible_cols:
                return possible_cols[0]
            # 尝试匹配其他别名
            for alias in aliases:
                if alias in df.columns:
                    return alias
                possible_cols = [col for col in df.columns if str(col).startswith(alias)]
                if possible_cols:
                    return possible_cols[0]

    # 方法3：前缀匹配（处理带#的列名）
    possible_cols = [col for col in df.columns if str(col).startswith(indicator)]
    if possible_cols:
        return possible_cols[0]

    # 方法4：去除空格后匹配
    indicator_no_space = indicator.replace(' ', '').replace('\u3000', '')
    for col in df.columns:
        col_no_space = str(col).replace(' ', '').replace('\u3000', '')
        if col_no_space == indicator_no_space:
            return col
        if col_no_space.startswith(indicator_no_space):
            return col

    # 方法5：部分匹配（宽松匹配）
    for col in df.columns:
        col_str = str(col)
        col_base = col_str.split('#')[0]  # 去除#后缀

        # 如果指标名是列名的子串
        if indicator in col_str or indicator in col_base:
            return col

        # 如果列名是指标名的子串
        if col_base in indicator:
            return col

    # 方法6：关键词匹配（最宽松）
    import re
    indicator_clean = re.sub(r'[（(].*?[）)]', '', indicator)  # 去除括号及内容
    indicator_clean = indicator_clean.strip()

    for col in df.columns:
        col_str = str(col).split('#')[0]  # 去除#后缀
        col_clean = re.sub(r'[（(].*?[）)]', '', col_str)
        col_clean = col_clean.strip()

        # 如果清理后的名称相同
        if indicator_clean == col_clean:
            return col

        # 如果指标名包含在列名中，或列名包含在指标名中
        if indicator_clean in col_clean or col_clean in indicator_clean:
            return col

    # 方法7：模糊匹配（允许1-2个字符不同）
    # 例如："平均红细胞血红浓度" vs "平均红细胞血红蛋白浓度"
    for col in df.columns:
        col_str = str(col).split('#')[0].strip()
        # 去除括号内容后比较
        col_clean = re.sub(r'[（(].*?[）)]', '', col_str).strip()
        indicator_clean_v2 = re.sub(r'[（(].*?[）)]', '', indicator).strip()

        # 如果长度相近（差距在3个字符以内）
        if abs(len(col_clean) - len(indicator_clean_v2)) <= 3:
            # 计算相似度：有多少个字符是相同的
            common_chars = sum(1 for c in indicator_clean_v2 if c in col_clean)
            similarity = common_chars / max(len(indicator_clean_v2), len(col_clean))

            # 如果相似度超过80%，认为匹配
            if similarity >= 0.8:
                return col

    return None

# ========== 图表生成函数 ==========

def plot_theme_table(athlete_df, theme_name, categories, ref_ranges, gender):
    """生成主题表格图 - 支持中英文双行显示"""
    if athlete_df.empty:
        return None, []

    latest_row = athlete_df.iloc[-1]
    latest_date = latest_row.get('DateStr', '未知')
    athlete_name = latest_row.get('Name', latest_row.get('Name_final', '未知'))

    cell_text = []
    cell_colors = []
    missing_indicators = []  # 记录缺失的指标

    # 状态中英文对照（包含优秀、良好等）
    status_translation = {
        '严重偏低': ('严重偏低', 'Severely Low'),
        '偏低': ('偏低', 'Low'),
        '正常': ('正常', 'Normal'),
        '良好': ('良好', 'Good'),
        '偏高': ('偏高', 'High'),
        '优秀': ('优秀', 'Excellent'),
        '严重偏高': ('严重偏高', 'Severely High'),
        '需注意': ('需注意', ' Attention'),
        '-': ('—', '—'),  # 无数据或未找到
        'N/A': ('—', '—'),  # 保留兼容
        '未找到': ('—', '—'),  # 保留兼容
    }

    # ⭐ 新增：格式化分类标题为多行显示
    def format_category_title(title):
        """
        将分类标题拆分为多行：
        输入："免疫防御（炎性监控）\nImmune Defense (Inflammatory Monitoring)"
        输出："免疫防御\nImmune Defense\n(Inflammatory Monitoring)"  ← 用英文括号！
        """
        lines = title.split('\n')
        if len(lines) != 2:
            return title  # 格式不对，保持原样
        
        cn_line = lines[0]  # 中文行：免疫防御（炎性监控）
        en_line = lines[1]  # 英文行：Immune Defense (Inflammatory Monitoring)
        
        # 提取中文主体（去掉中文括号部分）
        import re
        cn_match = re.match(r'(.+?)（.+?）', cn_line)
        if cn_match:
            cn_main = cn_match.group(1)  # 免疫防御
        else:
            cn_main = cn_line
        
        # 提取英文主体和括号内容
        en_match = re.match(r'(.+?)\s*\((.+?)\)', en_line)
        if en_match:
            en_main = en_match.group(1).strip()  # Immune Defense
            en_paren = en_match.group(2)  # Inflammatory Monitoring ← 用英文！
            # 有括号：3行，括号内容用英文
            return f"{cn_main}\n{en_main}\n({en_paren})"
        else:
            # 无括号：2行
            return f"{cn_main}\n{en_line}"

    for category_title, indicators in categories.items():
        # ⭐ 格式化分类标题为多行
        formatted_title = format_category_title(category_title)
        # 添加分类标题行（4列）- 居中对齐
        cell_text.append([formatted_title, '', '', ''])
        cell_colors.append([COLOR_CATEGORY_HEADER, COLOR_CATEGORY_HEADER, COLOR_CATEGORY_HEADER, COLOR_CATEGORY_HEADER])

        for col_key, name_tuple in indicators.items():
            # name_tuple是(中文名, 英文名)
            cn_name, en_name = name_tuple
            
            # 普通指标处理（包括睾酮/皮质醇比值）
            # 查找实际的列名
            actual_col = find_indicator_column(athlete_df, col_key)

            # 获取正常范围
            range_str = "—"
            if col_key in ref_ranges:
                ranges = ref_ranges[col_key]
                low_2 = ranges.get('low_2')
                high_2 = ranges.get('high_2')

                if pd.notna(low_2) and pd.notna(high_2):
                    # 两个值都存在，显示范围
                    range_str = f"{format_number(low_2)}-{format_number(high_2)}"
                elif pd.notna(low_2):
                    # 只有下限
                    range_str = f"≥{format_number(low_2)}"
                elif pd.notna(high_2):
                    # 只有上限
                    range_str = f"≤{format_number(high_2)}"

            if actual_col and actual_col in latest_row.index:
                val = latest_row[actual_col]
                
                if pd.notna(val):
                    # 🔧 处理带<或>符号的值
                    val_str_raw = str(val).strip()
                    
                    # ⭐ 特殊处理：保留<或>符号
                    if val_str_raw.startswith('<') or val_str_raw.startswith('>') or val_str_raw.startswith('＜') or val_str_raw.startswith('＞'):
                        val_str = val_str_raw  # 直接使用原始字符串
                        # 尝试提取数值进行评价
                        try:
                            num_str = val_str_raw.lstrip('<>＜＞').strip()
                            val_num = float(num_str)
                            status, bg_color, _ = get_indicator_status(col_key, val_num, ref_ranges, gender)
                        except:
                            status = "-"
                            bg_color = COLOR_NORMAL
                    else:
                        # 正常数值处理
                        try:
                            val = float(val)
                            # ⭐ 特殊处理：睾酮/皮质醇比值保留4位小数
                            if col_key == '睾酮/皮质醇比值':
                                val_str = f"{val:.4f}"
                            elif abs(val) >= 1000:
                                val_str = f"{val:.0f}"
                            elif abs(val) >= 100:
                                val_str = f"{val:.1f}"
                            else:
                                val_str = f"{val:.2f}"
                            status, bg_color, _ = get_indicator_status(col_key, val, ref_ranges, gender)
                        except (ValueError, TypeError):
                            val_str = "—"
                            status = "-"
                            bg_color = COLOR_NORMAL
                else:
                    val_str = "—"
                    status = "-"  # 无数据显示为"-"
                    bg_color = COLOR_NORMAL  # ⭐ 改为COLOR_NORMAL
            else:
                val_str = "—"
                status = "-"  # 未找到显示为"-"
                bg_color = COLOR_NORMAL  # ⭐ 改为COLOR_NORMAL
                missing_indicators.append((col_key, f"{cn_name}/{en_name}"))

            # 构建双行文本
            indicator_text = f"{cn_name}\n{en_name}"
            # 如果status是"-"，直接显示"-"，否则查询翻译
            if status == "-":
                status_text = "—"
            else:
                status_cn, status_en = status_translation.get(status, (status, status))
                status_text = f"{status_cn}\n{status_en}"
            
            cell_text.append([indicator_text, val_str, range_str, status_text])
            cell_colors.append([COLOR_NORMAL, bg_color, COLOR_NORMAL, bg_color])  # ⭐ 改为COLOR_NORMAL

    # 创建图表（4列，高清晰度）
    fig_height = len(cell_text) * 0.9 + 1.5  # 增加行高以容纳双行文本
    fig, ax = plt.subplots(figsize=(10, fig_height), dpi=150)
    ax.axis('off')

    # 列宽设置：保持原来的比例
    col_widths = [0.45, 0.18, 0.18, 0.19]  # 不改变列宽
    table = ax.table(
        cellText=cell_text,
        colLabels=['检测指标\nIndicator', '结果\nResult', '参考范围\nReference', '评价\nEvaluation'],
        cellColours=cell_colors,
        loc='center',
        cellLoc='center',
        colColours=[COLOR_TABLE_HEADER] * 4,
        colWidths=col_widths
    )

    table.auto_set_font_size(False)
    table.set_fontsize(FONTSIZE_VALUE)
    table.scale(1, TABLE_ROW_HEIGHT)  # 增加行高比例
    
    # 将分类标题颜色转换为RGB元组用于比较
    from matplotlib.colors import to_rgba
    category_color_rgba = to_rgba(COLOR_CATEGORY_HEADER)
    
    # 🔍 调试信息
    print(f"\n🎨 表格颜色调试信息:")
    print(f"   目标颜色: {COLOR_CATEGORY_HEADER}")
    print(f"   转换后: RGBA{category_color_rgba}")
    
    category_cell_count = 0  # 统计找到的分类标题单元格数量

    # 样式设置
    for (r, c), cell in table.get_celld().items():
        if r == 0:  # 表头
            cell.set_text_props(weight='bold', color='black', fontsize=FONTSIZE_HEADER)  # 黑色文字
            cell.set_edgecolor('#DDDDDD')  # 灰色边框
        else:
            # 获取当前单元格颜色（RGBA元组）
            cell_color = cell.get_facecolor()
            # 比较RGB值（忽略alpha通道）
            is_category = (abs(cell_color[0] - category_color_rgba[0]) < 0.01 and
                          abs(cell_color[1] - category_color_rgba[1]) < 0.01 and
                          abs(cell_color[2] - category_color_rgba[2]) < 0.01)
            
            if is_category:  # 分类标题
                category_cell_count += 1
                # ⭐ 分类标题：淡灰色边框（保留表格结构）
                cell.set_edgecolor('#DDDDDD')  # 淡灰边框
                cell.set_linewidth(1)          # 正常线宽
                
                if c == 0:  # 第一列：显示多行文字，居中
                    cell.set_text_props(weight='bold', color='black', ha='center', va='center', fontsize=FONTSIZE_CATEGORY)
                else:  # 其他列：隐藏文本
                    cell.set_text_props(visible=False)
                    # ⭐ 其他列的背景也设为分类标题颜色，确保整行一致
                    cell.set_facecolor(COLOR_CATEGORY_HEADER)
            else:  # 数据行
                cell.set_edgecolor('#DDDDDD')
                if r > 0 and c == 0:  # 指标名称列，左对齐
                    cell.set_text_props(ha='left', fontsize=FONTSIZE_INDICATOR)
                elif r > 0 and c in [1, 2]:  # 数值和范围列，较小字体
                    cell.set_text_props(fontsize=FONTSIZE_VALUE)
                elif r > 0 and c == 3:  # 评价列
                    cell.set_text_props(fontsize=FONTSIZE_STATUS)
    
    # 🔍 输出统计
    print(f"   找到分类标题单元格: {category_cell_count} 个")
    print(f"   预期数量: {len([k for k in cell_colors if k[0] == COLOR_CATEGORY_HEADER]) * 4}")

    # 获取中英文标题
    if theme_name in CATEGORY_NAMES:
        cn_title, en_title = CATEGORY_NAMES[theme_name]
        title_text = f"{athlete_name} ({gender}) - {cn_title}\n{en_title} ({latest_date})"
    else:
        theme_display = theme_name.split('_')[-1]
        title_text = f"{athlete_name} ({gender}) - {theme_display} ({latest_date})"
    
    # 一级标题：使用配置的字体大小，最小间距
    ax.set_title(title_text, fontsize=FONTSIZE_MAIN_TITLE, weight='bold', pad=2)

    plt.tight_layout()

    return fig, missing_indicators

def plot_trend_chart_multi(df, indicator, ref_ranges, selected_athletes, date_range, gender):
    """绘制多运动员对比趋势图"""

    # 查找实际的列名
    actual_col = find_indicator_column(df, indicator)
    if not actual_col:
        return None

    # 筛选日期范围
    if date_range and len(date_range) == 2:
        # 将date转换为datetime64以匹配df['Date']的类型
        start_date = pd.to_datetime(date_range[0])
        end_date = pd.to_datetime(date_range[1])
        df_filtered = df[(df['Date'] >= start_date) & (df['Date'] <= end_date)].copy()
    else:
        df_filtered = df.copy()

    if df_filtered.empty:
        return None

    # 获取名字列
    name_col = 'Name' if 'Name' in df_filtered.columns else 'Name_final'

    # 只保留有该指标数据的日期
    df_with_indicator = df_filtered[df_filtered[actual_col].notna()].copy()

    if df_with_indicator.empty:
        return None

    # 获取所有选中运动员中有数据的日期（去重排序）
    dates_with_data = set()
    for athlete in selected_athletes:
        athlete_data = df_with_indicator[df_with_indicator[name_col] == athlete]
        if not athlete_data.empty:
            dates_with_data.update(athlete_data['DateStr'].unique())

    # 如果没有任何数据，返回None
    if not dates_with_data:
        return None

    # 排序日期
    all_dates = sorted(list(dates_with_data))
    date_to_index = {date: i for i, date in enumerate(all_dates)}

    fig, ax = plt.subplots(figsize=(12, 7), dpi=150)
    ax.set_facecolor(COLOR_CHART_BG)

    # 协调配色列表（用于多运动员曲线）
    harmonious_colors = [
        '#4A90E2',  # 深海蓝
        '#D05A5E',  # 深砖红
        '#8BC1E9',  # 浅天蓝
        '#E89A9D',  # 浅柔红
        '#5C7CFA',  # 靛蓝
        '#9B59B6',  # 紫色
        '#1ABC9C',  # 青绿
        '#E67E22',  # 深橙
    ]

    # 确保有足够的颜色
    if len(selected_athletes) > len(harmonious_colors):
        colors = plt.cm.tab10(np.linspace(0, 1, len(selected_athletes)))
    else:
        colors = [harmonious_colors[i % len(harmonious_colors)] for i in range(len(selected_athletes))]

    # 先收集所有y值，用于确定范围
    all_y_values = []
    
    # 绘制每个运动员的数据
    for idx, (athlete, color) in enumerate(zip(selected_athletes, colors)):
        athlete_data = df_with_indicator[df_with_indicator[name_col] == athlete].copy()

        if athlete_data.empty:
            continue

        athlete_data = athlete_data.sort_values('Date')
        valid_data = athlete_data.dropna(subset=[actual_col])

        if len(valid_data) == 0:
            continue

        x_data = np.array([date_to_index[d] for d in valid_data['DateStr']])
        y_data = valid_data[actual_col].values
        all_y_values.extend(y_data)

        # 绘制平滑曲线
        if len(valid_data) > 1:
            try:
                x_smooth = np.linspace(x_data.min(), x_data.max(), 200)
                k = 2 if len(x_data) >= 3 else 1
                spl = make_interp_spline(x_data, y_data, k=k)
                y_smooth = spl(x_smooth)
                ax.plot(x_smooth, y_smooth, color=color, linewidth=2.5, label=athlete, alpha=0.8)
            except:
                ax.plot(x_data, y_data, color=color, linewidth=2.5, label=athlete, alpha=0.8)
        else:
            ax.plot(x_data, y_data, color=color, linewidth=2.5, label=athlete, linestyle='--', alpha=0.6)

        # 绘制数据点
        ax.plot(x_data, y_data, marker='o', markersize=8, markerfacecolor='white',
                markeredgecolor=color, markeredgewidth=2, linestyle='None')
        
        # 只为主运动员（第一个运动员）标注数据值
        if idx == 0:
            for x, y in zip(x_data, y_data):
                ax.text(x, y, f'{y:.1f}', 
                       fontsize=9, ha='center', va='bottom',
                       color=color, fontweight='bold',
                       bbox=dict(boxstyle='round,pad=0.3', facecolor='white', 
                                edgecolor=color, alpha=0.8, linewidth=1))

    # 在绘制数据后，添加理想范围标记
    if indicator in ref_ranges and len(all_y_values) > 0:
        ranges = ref_ranges[indicator]
        low_2 = ranges.get('low_2')
        high_2 = ranges.get('high_2')
        
        # 获取实际数据范围
        data_min = min(all_y_values)
        data_max = max(all_y_values)
        y_range = data_max - data_min
        
        # 情况1：同时有上下限（完整范围）
        if pd.notna(low_2) and pd.notna(high_2):
            ax.axhspan(low_2, high_2, color='#4A90E2', alpha=0.15, zorder=0, label='理想范围')
            ax.axhline(low_2, color=COLOR_SEVERE_LOW, linestyle=':', linewidth=1, alpha=0.7)
            ax.axhline(high_2, color=COLOR_SEVERE_HIGH, linestyle=':', linewidth=1, alpha=0.7)
        
        # 情况2：只有上限（如 < 300）
        elif pd.notna(high_2) and not pd.notna(low_2):
            # 从0或数据最小值往下一点开始
            y_min = min(0, data_min - y_range * 0.1)
            ax.axhspan(y_min, high_2, color='#4A90E2', alpha=0.15, zorder=0, label=f'理想范围 (< {high_2})')
            ax.axhline(high_2, color=COLOR_SEVERE_HIGH, linestyle=':', linewidth=1.5, alpha=0.7)
        
        # 情况3：只有下限（如 > 50）
        elif pd.notna(low_2) and not pd.notna(high_2):
            # 到数据最大值往上一点
            y_max = data_max + y_range * 0.1
            ax.axhspan(low_2, y_max, color='#4A90E2', alpha=0.15, zorder=0, label=f'理想范围 (> {low_2})')
            ax.axhline(low_2, color=COLOR_SEVERE_LOW, linestyle=':', linewidth=1.5, alpha=0.7)

    # 设置坐标轴 - 只显示有数据的日期
    ax.set_xticks(np.arange(len(all_dates)))
    ax.set_xticklabels(all_dates, rotation=45, ha='right')

    plt.title(f"{indicator} 趋势对比 ({gender})", fontsize=14, fontweight='bold')
    plt.xlabel('测试日期', fontsize=12)
    plt.ylabel(f'{indicator}', fontsize=12)
    plt.grid(axis='y', linestyle='--', alpha=0.5)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)

    # 图例
    plt.legend(loc='upper left', bbox_to_anchor=(1.01, 1), frameon=True)
    plt.tight_layout()

    return fig

def plot_radar_chart_with_baseline(athlete_df, radar_fields, lower_is_better, ref_ranges, athlete_name, baseline_athletes_df, gender):
    """
    绘制单个运动员的雷达图（最近4次测试）

    参数：
    - athlete_df: 主运动员的数据
    - radar_fields: 雷达图指标列表
    - lower_is_better: 逆指标列表
    - ref_ranges: 参考范围
    - athlete_name: 主运动员姓名
    - baseline_athletes_df: 用于计算baseline的所有运动员数据（包括主运动员）
    - gender: 性别
    """
    if athlete_df.empty:
        return None

    # 获取主运动员的最近4次数据
    last_4_dates = athlete_df['DateStr'].unique()[-4:]
    if len(last_4_dates) == 0:
        return None

    # 计算baseline统计值：使用对比运动员组的最近4次数据
    # 这样可以看到主运动员相对于对比组的表现
    baseline_stats = {}

    for field in radar_fields:
        actual_col = find_indicator_column(baseline_athletes_df, field)
        if actual_col:
            col_data = baseline_athletes_df[actual_col].dropna()
            if len(col_data) >= 2:
                baseline_stats[field] = {'mu': col_data.mean(), 'sigma': col_data.std()}
            else:
                baseline_stats[field] = {'mu': col_data.mean() if len(col_data) > 0 else 0, 'sigma': 1}
        else:
            baseline_stats[field] = {'mu': 0, 'sigma': 1}

    # 计算Z-score范围（用于设置坐标轴）
    athlete_z_scores = []
    for date in last_4_dates:
        date_row = athlete_df[athlete_df['DateStr'] == date]
        if date_row.empty:
            continue

        for field in radar_fields:
            actual_col = find_indicator_column(date_row, field)
            stats = baseline_stats.get(field)

            if not stats or stats['sigma'] == 0:
                z = 0
            else:
                if actual_col and actual_col in date_row.columns:
                    val = date_row[actual_col].values[0]
                    if pd.notna(val):
                        z = (val - stats['mu']) / stats['sigma']
                    else:
                        z = 0
                else:
                    z = 0

            if field in lower_is_better:
                z = -z
            athlete_z_scores.append(z)

    max_abs_z = max([abs(z) for z in athlete_z_scores]) if athlete_z_scores else 0
    limit = max(3, np.ceil(max_abs_z * 2) / 2)

    # 设置标签
    labels = [f + ('\n(逆)' if f in lower_is_better else '') for f in radar_fields]
    angles = np.linspace(0, 2 * np.pi, len(labels), endpoint=False).tolist()
    angles += angles[:1]

    # 创建图表（高清晰度）
    fig, ax = plt.subplots(figsize=(10, 10), subplot_kw=dict(polar=True), dpi=150)
    plt.ylim(-limit - 1.0, limit)

    # 绘制零线
    ax.plot(angles, [0] * len(angles), color='red', linewidth=2, linestyle='--', zorder=0.5)

    # ========== 绘制正常范围（浅绿色背景）==========
    normal_range_lower = []
    normal_range_upper = []
    
    for field in radar_fields:
        if field in ref_ranges:
            ranges = ref_ranges[field]
            low_2 = ranges.get('low_2')  # 正常范围下限
            high_2 = ranges.get('high_2')  # 正常范围上限
            stats = baseline_stats.get(field)
            
            if stats and stats['sigma'] != 0:
                # ⭐ 修复：分别处理下限和上限
                # 计算下限
                if pd.notna(low_2):
                    z_lower = (low_2 - stats['mu']) / stats['sigma']
                else:
                    z_lower = -2.5  # 没有下限，使用较小值
                
                # 计算上限
                if pd.notna(high_2):
                    z_upper = (high_2 - stats['mu']) / stats['sigma']
                else:
                    z_upper = 2.5  # 没有上限，使用较大值
                
                # 如果是逆指标，取反
                if field in lower_is_better:
                    z_lower, z_upper = -z_upper, -z_lower
                
                normal_range_lower.append(z_lower)
                normal_range_upper.append(z_upper)
            else:
                # 没有统计数据，使用默认值
                normal_range_lower.append(-1)
                normal_range_upper.append(1)
        else:
            # 没有参考范围，使用默认值
            normal_range_lower.append(-1)
            normal_range_upper.append(1)
    
    # 闭合多边形
    normal_range_lower.append(normal_range_lower[0])
    normal_range_upper.append(normal_range_upper[0])
    
    # 绘制理想范围区域（浅绿色填充）
    ax.fill_between(angles, normal_range_lower, normal_range_upper, 
                     color='#90EE90', alpha=0.2, zorder=1, label='理想范围')
    
    # 绘制理想范围边界线
    ax.plot(angles, normal_range_lower, color='#32CD32', linewidth=1.5, 
            linestyle=':', alpha=0.6, zorder=1)
    ax.plot(angles, normal_range_upper, color='#32CD32', linewidth=1.5, 
            linestyle=':', alpha=0.6, zorder=1)

    # 选择样式 - 最近4次测试
    styles = RADAR_STYLES[-len(last_4_dates):]

    # 绘制主运动员的最近4次数据
    for i, date in enumerate(last_4_dates):
        date_row = athlete_df[athlete_df['DateStr'] == date]
        if date_row.empty:
            continue

        values = []
        for field in radar_fields:
            actual_col = find_indicator_column(date_row, field)
            stats = baseline_stats.get(field)

            if not stats or stats['sigma'] == 0:
                z = 0
            else:
                if actual_col and actual_col in date_row.columns:
                    val = date_row[actual_col].values[0]
                    if pd.notna(val):
                        z = (val - stats['mu']) / stats['sigma']
                    else:
                        z = 0
                else:
                    z = 0

            if field in lower_is_better:
                z = -z
            values.append(z)

        values.append(values[0])
        style = styles[i]

        ax.plot(angles, values, color=style['color'], linewidth=style['linewidth'],
                linestyle=style['linestyle'], label=date, zorder=2)

        # 最新一次填充
        if i == len(last_4_dates) - 1:
            ax.fill(angles, values, color=style['color'], alpha=0.15, zorder=3)

    # 设置坐标轴
    ax.set_xticks(angles[:-1])
    ax.set_xticklabels(labels, size=11)

    # 数值刻度
    step = 1 if limit <= 3 else 2
    z_ticks = np.arange(-int(limit), int(limit) + 1, step)
    ax.set_yticks(z_ticks)
    ax.set_yticklabels([f'{i:.0f}' for i in z_ticks], color='grey', size=10)

    plt.title(f"{athlete_name} ({gender}) - 机能状态 Z-Score 雷达图",
              fontsize=16, y=1.08, fontweight='bold')
    plt.legend(loc='upper right', bbox_to_anchor=(1.3, 1.1))

    plt.tight_layout()

    return fig

# ========== 主应用 ==========


# ════════════════════════════════════════════════════════
# 【新增】问题与建议页 + PPT 导出 ── 合并自 blood_recs_addition.py
# ════════════════════════════════════════════════════════

# ─────────────────────────────────────────────────────────────────────────────
# 【新增函数 1】从 Excel「通用建议」Sheet 加载建议文本
# ─────────────────────────────────────────────────────────────────────────────
def load_advice_from_excel(file):
    """
    从参考范围 Excel 的「通用建议」Sheet 加载建议文本。
    返回 dict: { '指标名称': { 'severe_low': str, 'low': str, 'high': str, 'severe_high': str } }
    """
    try:
        df = pd.read_excel(file, sheet_name='通用建议')
        advice = {}
        col_map = {
            'severe_low':  '严重偏低 (<)',
            'low':         '偏低 (范围)',
            'normal':      '参考范围 (正常)',
            'high':        '偏高 (范围)',
            'severe_high': '严重偏高 (>)',
        }
        for _, row in df.iterrows():
            name = str(row.get('指标名称', '')).strip()
            if not name or name == 'nan':
                continue
            advice[name] = {}
            for key, col in col_map.items():
                val = row.get(col, '')
                if pd.notna(val) and str(val) != 'nan' and str(val).strip() not in ['-', '']:
                    advice[name][key] = str(val).strip()
                else:
                    advice[name][key] = ''
        return advice
    except Exception as e:
        return {}


# ─────────────────────────────────────────────────────────────────────────────
# 【新增函数 2】获取异常指标列表（含默认建议文本）
# ─────────────────────────────────────────────────────────────────────────────
def get_abnormal_indicators(athlete_df, ref_ranges, gender, advice_dict=None):
    """
    扫描最新一次测试，返回所有异常指标的列表（严重异常排前，轻度异常排后）。

    Returns list of dicts:
      { indicator, actual_col, value, value_str, status, bg_color, severity,
        default_advice, unit }
    """
    if athlete_df.empty:
        return []

    latest_row = athlete_df.iloc[-1]
    items = []

    for indicator, ranges in ref_ranges.items():
        actual_col = find_indicator_column(athlete_df, indicator)
        if not actual_col or actual_col not in latest_row.index:
            continue
        val = latest_row[actual_col]
        if pd.isna(val):
            continue
        try:
            val_float = float(val)
        except (ValueError, TypeError):
            continue

        status, bg_color, _ = get_indicator_status(indicator, val_float, ref_ranges, gender)
        if status not in ['严重偏低', '严重偏高', '偏低', '偏高']:
            continue

        severity = 'severe' if '严重' in status else 'slight'

        # 格式化数值
        if indicator == '睾酮/皮质醇比值':
            val_str = f"{val_float:.4f}"
        elif abs(val_float) >= 1000:
            val_str = f"{val_float:.0f}"
        elif abs(val_float) >= 100:
            val_str = f"{val_float:.1f}"
        else:
            val_str = f"{val_float:.2f}"

        # 获取默认建议文本
        default_advice = ''
        if advice_dict and indicator in advice_dict:
            key_map = {'严重偏低': 'severe_low', '偏低': 'low',
                       '偏高': 'high', '严重偏高': 'severe_high'}
            adv = advice_dict[indicator].get(key_map.get(status, ''), '')
            if adv:
                default_advice = adv[:200]   # 最多200字

        # 推断单位（从 THEME_CONFIG 中找）
        unit = _get_indicator_unit(indicator, athlete_df, actual_col)

        items.append({
            'indicator': indicator,
            'actual_col': actual_col,
            'value': val_float,
            'value_str': val_str,
            'status': status,
            'bg_color': bg_color,
            'severity': severity,
            'default_advice': default_advice,
            'unit': unit,
        })

    # 严重在前，轻度在后，同等级按指标名排序
    items.sort(key=lambda x: (0 if x['severity'] == 'severe' else 1, x['indicator']))
    return items


def _get_indicator_unit(indicator, df, actual_col):
    """尝试推断指标单位（可扩展）"""
    unit_hints = {
        '血红蛋白': 'g/L', '铁蛋白': 'ng/ml', '睾酮': 'ng/dL', '游离睾酮': 'ng/dL',
        '皮质醇': 'μg/dL', '肌酸激酶': 'U/L', '血尿素': 'mmol/L', '尿酸': 'μmol/L',
        '血糖': 'mmol/L', '维生素D': 'ng/ml', '维生素D3': 'ng/ml', '维生素B12': 'pg/ml',
        '叶酸': 'μg/L', '白细胞': '10⁹/L', '超敏C反应蛋白': 'mg/L',
        '钾': 'mmol/L', '钠': 'mmol/L', '氯': 'mmol/L', '渗透压': 'mOsm/L',
    }
    return unit_hints.get(indicator, '')


# ─────────────────────────────────────────────────────────────────────────────
# 【新增函数 3】生成「问题与建议」matplotlib 图
# ─────────────────────────────────────────────────────────────────────────────
def plot_recommendations_page(athlete_name, gender, latest_date, items_with_advice):
    """
    生成「问题与建议」页的 matplotlib 图表（可直接插入 PPT）。

    Parameters
    ----------
    athlete_name : str
    gender       : str
    latest_date  : str  最新测试日期
    items_with_advice : list of dicts
        每个 dict 需含 keys:
        indicator, value_str, status, severity, unit, advice (可编辑后传入)
    """
    if not items_with_advice:
        items_with_advice = [{
            'indicator': '综合评估',
            'value_str': '—',
            'status': '正常',
            'severity': 'ok',
            'unit': '',
            'advice': '各项指标均在正常范围内，继续保持当前训练与营养方案，定期监测指标变化。',
        }]

    n = len(items_with_advice)
    ITEM_H  = 1.3      # 每行高度（英寸）
    HDR_H   = 0.70     # 标题栏高度
    INFO_H  = 0.55     # 运动员信息栏高度
    EXTRA_H = 0.15     # 底部留白
    fig_h   = HDR_H + INFO_H + n * ITEM_H + EXTRA_H

    fig = plt.figure(figsize=(16, fig_h), dpi=130)
    fig.patch.set_facecolor('white')
    ax = fig.add_axes([0, 0, 1, 1])
    ax.axis('off')
    ax.set_xlim(0, 16)
    ax.set_ylim(0, fig_h)

    # ── 标题栏 ──
    y = fig_h
    ax.add_patch(mpatches.FancyBboxPatch((0, y - HDR_H), 16, HDR_H,
        boxstyle="square,pad=0", facecolor='#1a3a5c', edgecolor='none', zorder=1))
    ax.text(8, y - HDR_H / 2, '重要启示 · Key Insights',
            ha='center', va='center', fontsize=20, fontweight='bold',
            color='white', zorder=2)

    # ── 运动员信息栏 ──
    y -= HDR_H
    ax.add_patch(mpatches.FancyBboxPatch((0, y - INFO_H), 16, INFO_H,
        boxstyle="square,pad=0", facecolor='#EEF2F8', edgecolor='none', zorder=1))
    ax.text(0.35, y - INFO_H / 2,
            f'运动员：{athlete_name}  ({gender})      最新测试日期：{latest_date}',
            ha='left', va='center', fontsize=12.5, fontweight='bold',
            color='#1a3a5c', zorder=2)

    # 统计徽章
    sev_n = sum(1 for i in items_with_advice if i['severity'] == 'severe')
    sli_n = sum(1 for i in items_with_advice if i['severity'] == 'slight')
    badge_x = 15.7
    if sev_n:
        ax.text(badge_x, y - INFO_H / 2 - 0.01,
                f'⬤ 严重异常 {sev_n} 项',
                ha='right', va='center', fontsize=11, fontweight='bold',
                color='#c62828', zorder=2)
        badge_x -= 2.8
    if sli_n:
        ax.text(badge_x, y - INFO_H / 2 - 0.01,
                f'⬤ 轻度异常 {sli_n} 项',
                ha='right', va='center', fontsize=11, fontweight='bold',
                color='#e65100', zorder=2)

    # 列标题
    ax.text(1.05, y - INFO_H / 2, '指标', ha='left', va='center',
            fontsize=9.5, color='#888888', zorder=2)
    ax.text(5.70, y - INFO_H / 2, '建议措施 / Recommendations', ha='left', va='center',
            fontsize=9.5, color='#888888', zorder=2)
    y -= INFO_H

    # ── 各条目 ──
    def wrap_cn(text, chars=52):
        """简单中文换行"""
        if not text or str(text) == 'nan':
            return []
        text = str(text)
        return [text[i:i+chars] for i in range(0, len(text), chars)]

    for idx, item in enumerate(items_with_advice):
        row_bg = '#FFFFFF' if idx % 2 == 0 else '#F8FAFD'
        ax.add_patch(mpatches.FancyBboxPatch((0, y - ITEM_H), 16, ITEM_H,
            boxstyle="square,pad=0", facecolor=row_bg, edgecolor='none', zorder=1))

        # 颜色方案
        if item['severity'] == 'severe':
            c_bg, c_num, c_stat = '#FFCDD2', '#B71C1C', '#c62828'
        elif item['severity'] == 'slight':
            c_bg, c_num, c_stat = '#FFF9C4', '#E65100', '#e65100'
        else:
            c_bg, c_num, c_stat = '#C8E6C9', '#2E7D32', '#2e7d32'

        # 序号圆圈
        circle = plt.Circle((0.45, y - ITEM_H / 2), 0.22,
                             color=c_bg, zorder=3)
        ax.add_patch(circle)
        ax.text(0.45, y - ITEM_H / 2, str(idx + 1),
                ha='center', va='center', fontsize=12,
                fontweight='bold', color=c_num, zorder=4)

        # 指标名 + 评价
        ax.text(0.90, y - ITEM_H * 0.35,
                item['indicator'],
                ha='left', va='center', fontsize=13,
                fontweight='bold', color='#1a1a1a', zorder=2)
        unit_str = f" ({item['unit']})" if item.get('unit') else ''
        ax.text(0.90, y - ITEM_H * 0.65,
                f"{item['status']}  当前值: {item['value_str']}{unit_str}",
                ha='left', va='center', fontsize=10.5,
                color=c_stat, zorder=2)

        # 竖分隔线
        ax.plot([5.55, 5.55], [y - ITEM_H + 0.08, y - 0.08],
                color='#DDDDDD', linewidth=0.8, zorder=2)

        # 建议文本（自动换行，最多4行）
        advice = item.get('advice', '') or ''
        adv_lines = wrap_cn(advice, chars=55)
        if adv_lines:
            LINE_H = 0.29
            total_th = min(len(adv_lines), 4) * LINE_H
            start_y = y - ITEM_H / 2 + total_th / 2 - LINE_H * 0.5
            for li, line in enumerate(adv_lines[:4]):
                ax.text(5.80, start_y - li * LINE_H,
                        line, ha='left', va='center',
                        fontsize=10.5, color='#333333', zorder=2)
        else:
            ax.text(5.80, y - ITEM_H / 2,
                    '请咨询队医获取详细干预方案',
                    ha='left', va='center', fontsize=10.5,
                    color='#AAAAAA', zorder=2)

        # 底部分隔线
        ax.plot([0, 16], [y - ITEM_H, y - ITEM_H],
                color='#E5EBF3', linewidth=0.5, zorder=2)
        y -= ITEM_H

    plt.tight_layout(pad=0)
    return fig


# ─────────────────────────────────────────────────────────────────────────────
# 【新增函数 4】将多个 matplotlib 图导出为 PPTX 字节流
# ─────────────────────────────────────────────────────────────────────────────
def export_figures_to_pptx(figures_with_titles, slide_w_inch=13.33, slide_h_inch=7.5):
    """
    将 matplotlib 图列表导出为 PPTX。
    参数:
        figures_with_titles : list of (fig, title_str)
        slide_w_inch        : 幻灯片宽度（英寸），默认 13.33（WIDE）
        slide_h_inch        : 幻灯片高度（英寸），默认 7.5
    返回:
        bytes  (PPTX 二进制内容) 或 None（失败时）
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        prs = Presentation()
        prs.slide_width  = Inches(slide_w_inch)
        prs.slide_height = Inches(slide_h_inch)
        blank_layout = prs.slide_layouts[6]   # 完全空白的版式

        for fig, title in figures_with_titles:
            if fig is None:
                continue

            slide = prs.slides.add_slide(blank_layout)

            # 将 matplotlib 图渲染为 PNG（内存中）
            img_buf = BytesIO()
            fig.savefig(img_buf, format='png', dpi=150,
                        bbox_inches='tight',
                        facecolor='white', edgecolor='none')
            img_buf.seek(0)

            # 全页插图
            slide.shapes.add_picture(
                img_buf,
                Inches(0), Inches(0),
                Inches(slide_w_inch), Inches(slide_h_inch)
            )

            # 在备注栏写入标题（方便后续整理）
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = title

        pptx_buf = BytesIO()
        prs.save(pptx_buf)
        return pptx_buf.getvalue()

    except ImportError:
        st.error("❌ 缺少 python-pptx 库，请运行：pip install python-pptx")
        return None
    except Exception as e:
        st.error(f"❌ PPT 导出失败：{e}")
        import traceback
        st.error(traceback.format_exc())
        return None


# ─────────────────────────────────────────────────────────────────────────────
# 【main() 修改说明】
# 在原 app__11_.py 的 main() 中做如下改动：
# ─────────────────────────────────────────────────────────────────────────────

def _main_changes_description():
    """
    以下是对原 main() 函数需要做的改动，逐处说明。
    （此函数仅作说明，不需要执行）

    1. 在 use_custom_ranges 代码块里，加载建议文本：
       ─────────────────────────────────────────────
       if use_custom_ranges and custom_ranges_file is not None:
           ...（原有代码）...
           # 【新增】同时加载建议文本
           advice_dict = load_advice_from_excel(custom_ranges_file)
           st.session_state['advice_dict'] = advice_dict
       else:
           st.session_state.setdefault('advice_dict', {})

    2. 在 tab 定义处改为 5 个选项卡：
       ─────────────────────────────────────────────
       # 原代码：
       # 初始化 PPT 图表存储
    if 'ppt_figs' not in st.session_state:
        st.session_state['ppt_figs'] = {}

    tab1, tab2, tab3, tab4, tab5 = st.tabs([
        "📋 主题表格", "📈 趋势对比", "🎯 雷达图", "📊 数据表", "📝 问题与建议 & PPT"
    ])
       # 改为：
       tab1, tab2, tab3, tab4, tab5 = st.tabs(
           ["📋 主题表格", "📈 趋势对比", "🎯 雷达图", "📊 数据表", "📝 问题与建议 & PPT"]
       )
       # 同时在 tabs 之前初始化 session_state：
       if 'ppt_figs' not in st.session_state:
           st.session_state['ppt_figs'] = {}

    3. 在 tab1「生成主题表格」按钮的 with st.spinner 里，把生成的图存进 session_state：
       ─────────────────────────────────────────────
       # 原代码（找到这段循环）：
       for theme_name, categories in THEME_CONFIG.items():
           ...
           if result:
               fig, missing = result
               if fig:
                   st.pyplot(fig)
                   plt.close()
       # 改为（仅加一行 session_state 存储，close 前）：
       for theme_name, categories in THEME_CONFIG.items():
           ...
           if result:
               fig, missing = result
               if fig:
                   st.pyplot(fig)
                   st.session_state['ppt_figs'][theme_name] = fig   # ← 新增
                   plt.close(fig)

    4. 在 tab3「生成雷达图」按钮的 if fig: 里，同样存图：
       ─────────────────────────────────────────────
       if fig:
           st.pyplot(fig)
           st.session_state['ppt_figs']['radar'] = fig              # ← 新增
           plt.close(fig)
           ...

    5. 将 tab5 的完整内容（见下方 tab5_content() 函数）加入 with tab5: 块。
    """
    pass


# ─────────────────────────────────────────────────────────────────────────────
# 【Tab5 完整实现】
# 在 main() 里加入： with tab5:  →  调用本函数内容
# ─────────────────────────────────────────────────────────────────────────────
def tab5_content(athlete_name, gender, athlete_df, ref_ranges):
    """
    Tab5「📝 问题与建议 & PPT」的完整 Streamlit 内容。
    直接把本函数 body 粘贴到 `with tab5:` 块中即可。
    """
    st.subheader("📝 问题与建议页（可编辑 → 插入 PPT）")

    advice_dict = st.session_state.get('advice_dict', {})

    # ── 获取异常指标 ──────────────────────────────────────────
    items = get_abnormal_indicators(athlete_df, ref_ranges, gender, advice_dict)
    latest_date = athlete_df.iloc[-1].get('DateStr', '未知') if not athlete_df.empty else '未知'

    if not items:
        st.success("✅ 所有指标均在参考范围内，无异常。")
    else:
        sev_n = sum(1 for i in items if i['severity'] == 'severe')
        sli_n = sum(1 for i in items if i['severity'] == 'slight')
        c1, c2, c3 = st.columns(3)
        c1.metric("运动员", athlete_name)
        c2.metric("🔴 严重异常", f"{sev_n} 项", delta=None)
        c3.metric("🟡 轻度异常", f"{sli_n} 项", delta=None)

    st.markdown("---")
    st.markdown("#### 逐项编辑建议措施")
    st.caption("每个异常指标的建议文本均可自由修改，修改后点击「生成预览」查看效果。")

    # ── 可编辑文本框 ──────────────────────────────────────────
    edited_items = []
    for item in items:
        sev_icon = "🔴" if item['severity'] == 'severe' else "🟡"
        with st.expander(
            f"{sev_icon} **{item['indicator']}** — {item['status']}  "
            f"（当前值: {item['value_str']} {item.get('unit','')}）",
            expanded=(item['severity'] == 'severe')
        ):
            adv_key = f"adv_{athlete_name}_{item['indicator']}"
            edited_adv = st.text_area(
                "建议措施（可自由修改）",
                value=item['default_advice'],
                key=adv_key,
                height=100,
                help="此文本将显示在问题与建议页右侧。"
            )
            edited_items.append({**item, 'advice': edited_adv})

    # 若无异常指标，仍允许写一段综合评估
    if not items:
        adv_ok_key = f"adv_{athlete_name}_综合评估"
        ok_adv = st.text_area(
            "综合评估建议",
            value="各项指标均在正常范围内，继续保持当前训练与营养方案，定期监测指标变化。",
            key=adv_ok_key,
            height=80
        )
        edited_items = [{
            'indicator': '综合评估',
            'value_str': '—', 'status': '正常',
            'severity': 'ok', 'unit': '',
            'advice': ok_adv,
        }]

    st.markdown("---")

    # ── 预览按钮 ──────────────────────────────────────────────
    col_prev, col_export = st.columns([1, 1])

    with col_prev:
        if st.button("🔄 生成问题与建议预览", use_container_width=True):
            with st.spinner("正在生成预览…"):
                fig_recs = plot_recommendations_page(
                    athlete_name, gender, latest_date, edited_items
                )
                st.session_state['ppt_figs']['recs'] = fig_recs
                st.pyplot(fig_recs)
                st.success("✅ 预览生成完成！已保存，可直接导出 PPT。")
                plt.close(fig_recs)

    # ── PPT 导出 ──────────────────────────────────────────────
    st.markdown("---")
    st.subheader("📤 导出 PPT")

    ppt_figs = st.session_state.get('ppt_figs', {})
    available = [k for k, v in ppt_figs.items() if v is not None]

    if not available:
        st.info("ℹ️ 请先在各选项卡中生成图表（主题表格、雷达图、问题与建议），再导出 PPT。")
    else:
        # 显示已生成的图表列表，供用户确认
        st.markdown("**以下图表将被导出（已生成的图表）：**")
        label_map = {
            '1_调控与指挥中心': '📊 System I: 调控与指挥中心',
            '2_执行与代谢系统': '📊 System II: 执行与代谢系统',
            '3_循环与运载系统': '📊 System III: 循环与运载系统',
            '4_后勤保障与维护': '📊 System IV: 后勤保障与维护',
            'radar': '🕸️ 雷达图',
            'recs':  '💡 问题与建议',
        }
        # 趋势图
        for k in available:
            if k.startswith('trend_'):
                label_map[k] = f'📈 趋势图: {k.replace("trend_","")}'

        sel_keys = st.multiselect(
            "选择要导入 PPT 的图表",
            options=available,
            default=available,
            format_func=lambda k: label_map.get(k, k),
        )

        # 幻灯片顺序设置
        st.markdown("**幻灯片顺序**（从上到下依次排列）")
        order_map = {
            '1_调控与指挥中心': 0, '2_执行与代谢系统': 1,
            '3_循环与运载系统': 2, '4_后勤保障与维护': 3,
            'radar': 4, 'recs': 99,
        }
        sel_sorted = sorted(sel_keys,
                            key=lambda k: order_map.get(k, 50) if not k.startswith('trend_')
                            else 10 + int(k.replace('trend_', '') or 0))

        for i, k in enumerate(sel_sorted):
            st.caption(f"  第 {i+1} 页 — {label_map.get(k, k)}")

        if st.button("🚀 生成并下载 PPT", type="primary", use_container_width=True):
            with st.spinner("正在生成 PPT…"):
                # 按顺序收集图表
                figs_to_export = []
                for k in sel_sorted:
                    fig = ppt_figs.get(k)
                    if fig is not None:
                        figs_to_export.append((fig, label_map.get(k, k)))

                if not figs_to_export:
                    st.warning("⚠️ 没有可导出的图表。")
                else:
                    pptx_bytes = export_figures_to_pptx(figs_to_export)
                    if pptx_bytes:
                        fname = f"{athlete_name}_血液生化报告.pptx"
                        st.download_button(
                            label=f"📥 下载 {fname}",
                            data=pptx_bytes,
                            file_name=fname,
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            use_container_width=True,
                        )
                        st.success(f"✅ PPT 生成成功！共 {len(figs_to_export)} 张幻灯片。")



def main():
    st.title("🏃 运动员血液指标分析系统")
    st.markdown("**包含：表格图、多运动员趋势对比、雷达图**")
    st.markdown("---")

    # === 侧边栏 ===
    st.sidebar.header("📂 数据上传")

    # 数据文件上传
    uploaded_file = st.sidebar.file_uploader(
        "1️⃣ 上传血液数据Excel",
        type=['xlsx', 'xls'],
        help="请上传包含'月周测试指标'工作表的Excel文件",
        key="data_file"
    )

    # 参考范围文件上传
    st.sidebar.markdown("---")
    st.sidebar.markdown("**📊 参考范围设置**")

    use_custom_ranges = st.sidebar.checkbox(
        "使用自定义参考范围",
        value=False,
        help="勾选后可上传自己的参考范围Excel文件"
    )

    custom_ranges_file = None
    if use_custom_ranges:
        custom_ranges_file = st.sidebar.file_uploader(
            "2️⃣ 上传参考范围Excel",
            type=['xlsx', 'xls'],
            help="Excel文件需包含'参考范围'工作表",
            key="ranges_file"
        )

    if uploaded_file is None:
        st.info("👈 请在左侧上传Excel数据文件")
        st.stop()

    # === 加载参考范围 ===
    if use_custom_ranges and custom_ranges_file is not None:
        with st.spinner("正在加载自定义参考范围..."):
            male_ref_ranges, female_ref_ranges = load_reference_ranges_from_excel(custom_ranges_file)
            if male_ref_ranges and female_ref_ranges:
                st.sidebar.success(f"✅ 已加载自定义范围（男:{len(male_ref_ranges)}项，女:{len(female_ref_ranges)}项）")
            else:
                st.sidebar.warning("⚠️ 自定义范围加载失败，使用默认范围")
                male_ref_ranges = MALE_REF_RANGES
                female_ref_ranges = FEMALE_REF_RANGES
    else:
        # 使用默认范围
        male_ref_ranges = MALE_REF_RANGES
        female_ref_ranges = FEMALE_REF_RANGES
        if not use_custom_ranges:
            st.sidebar.info("ℹ️ 使用默认参考范围")

    # === 数据加载 ===
    with st.spinner("⏳ 正在读取数据（大文件首次加载约30秒，请稍候）..."):
        df = load_data_multisheet(uploaded_file)

        if df is None:
            st.stop()

        df = clean_data_final(df)

        if df is None or len(df) == 0:
            st.error("❌ 数据清洗后为空")
            st.stop()

    st.success(f"🎉 数据准备完成：共 {len(df)} 条记录")

    # === 数据预览 ===
    with st.expander("👀 查看数据"):
        st.write("**前20行：**")
        st.write(df.head(20))

    st.markdown("---")

    # === 用户选择 ===
    col1, col2 = st.columns(2)

    with col1:
        gender = st.selectbox("选择性别", ["男", "女"])

    # 筛选性别
    if '性别' in df.columns:
        gender_df = df[df['性别'] == gender].copy()
    else:
        st.warning("⚠️ 数据中没有'性别'列")
        gender_df = df.copy()

    if len(gender_df) == 0:
        st.warning(f"⚠️ 没有{gender}运动员的数据")
        st.stop()

    # 获取运动员列表
    name_col = None
    for possible_name in ['Name_final', 'Name', '姓名']:
        if possible_name in gender_df.columns:
            name_col = possible_name
            break

    if not name_col:
        st.error("❌ 未找到姓名列")
        st.stop()

    athletes = sorted(gender_df[name_col].dropna().unique())
    ref_ranges = male_ref_ranges if gender == "男" else female_ref_ranges

    with col2:
        athlete_name = st.selectbox(
            "选择运动员",
            athletes,
            help=f"共 {len(athletes)} 名{gender}运动员"
        )

    # 筛选运动员数据
    athlete_df = gender_df[gender_df[name_col] == athlete_name].copy()

    date_col = 'Date' if 'Date' in athlete_df.columns else 'Date_auto'
    if date_col in athlete_df.columns:
        athlete_df = athlete_df.sort_values(date_col)

    st.info(f"📊 **{athlete_name}**（{gender}）- 共 {len(athlete_df)} 次测试")

    st.markdown("---")

    # === 获取所有可用的数值指标 ===
    # 排除非指标列
    exclude_cols = ['Name', 'Name_final', '姓名', 'Date', 'Date_auto', '日期', 'DateStr', 
                    '性别', 'Gender', '编号', 'ID', 'Unnamed: 0']
    all_numeric_indicators = []
    for col in gender_df.columns:
        if col not in exclude_cols:
            # 检查是否是数值列
            try:
                if gender_df[col].dtype in ['float64', 'int64'] or pd.to_numeric(gender_df[col], errors='coerce').notna().any():
                    all_numeric_indicators.append(col)
            except:
                pass
    
    # 如果没有找到数值列，使用默认的TREND_INDICATORS
    if not all_numeric_indicators:
        all_numeric_indicators = TREND_INDICATORS

    # === 功能选项卡 ===
    tab1, tab2, tab3, tab4 = st.tabs(["📋 主题表格", "📈 趋势对比", "🎯 雷达图", "📊 数据表"])

    # --- Tab 1: 主题表格 ---
    with tab1:
        st.subheader("最新数据主题表格")
        st.markdown("显示最新一次测试的各项指标，使用五档判断")

        if st.button("🚀 生成主题表格", type="primary", use_container_width=True):
            with st.spinner("正在生成表格..."):

                for theme_name, categories in THEME_CONFIG.items():
                    st.markdown(f"<h2 style='margin-bottom: {TITLE_TABLE_SPACING}em;font-size: {FONTSIZE_MAIN_TITLE}px;'>{theme_name.split('_')[-1]}</h3>", unsafe_allow_html=True)
                    result = plot_theme_table(athlete_df, theme_name, categories, ref_ranges, gender)

                    if result:
                        fig, missing = result
                        if fig:
                            st.pyplot(fig)
                            st.session_state['ppt_figs'][theme_name] = fig
                            plt.close(fig)
                        else:
                            st.info(f"ℹ️ {theme_name} 数据不足")
                    else:
                        st.info(f"ℹ️ {theme_name} 数据不足")

                st.success("✅ 表格生成完成！")

    # --- Tab 2: 趋势对比 ---
    with tab2:
        st.subheader("多运动员趋势对比")
        st.markdown("可以选择多个运动员和日期范围进行对比")

        # 选择对比运动员
        compare_athletes = st.multiselect(
            "选择对比运动员（可多选）",
            athletes,
            default=[athlete_name],
            help="选择要对比的运动员"
        )

        # 日期范围选择
        if date_col in gender_df.columns:
            min_date = gender_df[date_col].min()
            max_date = gender_df[date_col].max()

            date_range = st.date_input(
                "选择日期范围",
                value=(min_date, max_date),
                min_value=min_date,
                max_value=max_date,
                help="选择要分析的日期范围"
            )
        else:
            date_range = None

        # 选择指标
        # 构建默认选择：优先使用TREND_INDICATORS中存在于数据的指标
        default_trend = [ind for ind in TREND_INDICATORS if ind in all_numeric_indicators]
        if not default_trend and all_numeric_indicators:
            # 如果TREND_INDICATORS中的指标都不在数据中，使用前3个
            default_trend = all_numeric_indicators[:3] if len(all_numeric_indicators) >= 3 else all_numeric_indicators
        
        selected_indicators = st.multiselect(
            "选择要分析的指标",
            all_numeric_indicators,
            default=default_trend,
            help="选择要绘制趋势图的指标（可选择所有数值指标）"
        )

        if st.button("🚀 生成趋势对比图", type="primary", use_container_width=True):
            if not compare_athletes:
                st.warning("⚠️ 请至少选择一个运动员")
            elif not selected_indicators:
                st.warning("⚠️ 请至少选择一个指标")
            else:
                with st.spinner("正在生成趋势图..."):
                    for indicator in selected_indicators:
                        st.markdown(f"### {indicator}")
                        fig = plot_trend_chart_multi(
                            gender_df, indicator, ref_ranges,
                            compare_athletes, date_range, gender
                        )
                        if fig:
                            st.pyplot(fig)
                            plt.close()
                        else:
                            st.info(f"ℹ️ {indicator} 数据不足")

                    st.success("✅ 趋势图生成完成！")

    # --- Tab 3: 雷达图 ---
    with tab3:
        st.subheader(f"{athlete_name}的机能状态雷达图")
        st.markdown(f"显示**{athlete_name}**最近4次测试的Z-Score雷达图")

        # 说明Z-Score计算方式
        st.info("💡 **Z-Score计算说明**：使用对比运动员组的数据作为基准，计算该运动员相对于组内的表现")

        # 选择对比运动员组（用于计算baseline）
        radar_athletes = st.multiselect(
            "选择对比运动员组（用于计算Z-Score基准）",
            athletes,
            default=[athlete_name],
            help="选择的运动员将作为基准组，用于计算Z-Score的均值和标准差",
            key="radar_athletes"
        )

        # 选择雷达图指标
        # 构建默认选择：优先使用RADAR_FIELDS中存在于数据的指标
        default_radar = [ind for ind in RADAR_FIELDS if ind in all_numeric_indicators]
        if not default_radar and all_numeric_indicators:
            # 如果RADAR_FIELDS中的指标都不在数据中，使用前8个
            default_radar = all_numeric_indicators[:8] if len(all_numeric_indicators) >= 8 else all_numeric_indicators
        
        radar_indicators = st.multiselect(
            "选择雷达图指标",
            all_numeric_indicators,
            default=default_radar,
            help="选择要在雷达图中显示的指标（建议4-10个，可选择所有数值指标）"
        )

        # 选择逆指标（值越低越好的指标）
        st.markdown("**逆指标设置**（值越低越好的指标）")
        lower_better = st.multiselect(
            "选择逆指标",
            radar_indicators,
            default=[ind for ind in LOWER_IS_BETTER if ind in radar_indicators],
            help="这些指标在雷达图中会取反（如肌酸激酶、血尿素等）"
        )

        if st.button("🚀 生成雷达图", type="primary", use_container_width=True, key="radar_btn"):
            if not radar_athletes:
                st.warning("⚠️ 请至少选择一个对比运动员")
            elif not radar_indicators:
                st.warning("⚠️ 请至少选择一个指标")
            elif len(radar_indicators) < 3:
                st.warning("⚠️ 请至少选择3个指标，雷达图效果更好")
            else:
                with st.spinner("正在生成雷达图..."):
                    # 获取对比运动员组的最近4次数据（用于计算baseline）
                    baseline_data_list = []
                    for comp_athlete in radar_athletes:
                        comp_athlete_df = gender_df[gender_df[name_col] == comp_athlete].sort_values('Date')
                        if not comp_athlete_df.empty:
                            # 获取该运动员的最近4次数据
                            last_4 = comp_athlete_df.tail(4)
                            baseline_data_list.append(last_4)

                    if baseline_data_list:
                        baseline_df = pd.concat(baseline_data_list, ignore_index=True)

                        # 生成雷达图：只画主运动员的近4次，但用baseline_df计算Z值
                        fig = plot_radar_chart_with_baseline(
                            athlete_df, radar_indicators, lower_better,
                            ref_ranges, athlete_name, baseline_df, gender
                        )

                        if fig:
                            st.pyplot(fig)
                            st.session_state['ppt_figs']['radar'] = fig
                            plt.close(fig)
                            st.success("✅ 雷达图生成完成！")

                            # 添加说明
                            st.markdown("---")
                            st.markdown("### 📖 雷达图说明")
                            st.markdown(f"""
                            - **显示内容**：{athlete_name}的最近4次测试
                            - **对比基准**：使用{len(radar_athletes)}个运动员的最近4次数据计算均值和标准差
                            - **Z-Score含义**：
                              - **0**：等于基准组平均水平
                              - **正值**：高于基准组平均水平
                              - **负值**：低于基准组平均水平
                            - **逆指标**：标记"(逆)"的指标已取反显示（值越低越好）
                            - **线条样式**：
                              - 蓝色虚点线：第1次测试
                              - 橙色点划线：第2次测试
                              - 绿色虚线：第3次测试
                              - 红色实线+填充：第4次测试（最新）
                            - **解读要点**：图形越向外，表现越好；图形越规则，机能越均衡
                            """)
                        else:
                            st.info("ℹ️ 数据不足，无法生成雷达图")
                    else:
                        st.warning("⚠️ 对比运动员组没有足够的数据")

    # --- Tab 4: 数据表 ---
    with tab4:
        st.subheader("完整数据表")
        st.write(athlete_df)

        try:
            csv = athlete_df.to_csv(index=False, encoding='utf-8-sig')
            st.download_button(
                label="📥 下载CSV数据",
                data=csv,
                file_name=f"{athlete_name}_数据.csv",
                mime="text/csv"
            )
        except:
            st.warning("CSV下载功能暂时不可用")

    # --- Tab 5: 问题与建议 & PPT 导出 ---
    with tab5:
        tab5_content(athlete_name, gender, athlete_df, ref_ranges)

if __name__ == "__main__":
    main()
